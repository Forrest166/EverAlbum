#!/usr/bin/env python3
"""
📷 智能照片相册生成器 v4.0

v4.0 新增:
  1. 事件页直接显示 EXIF 经纬度，不再调用网络地理编码
  2. PPTX 采用 A4 尺寸，与 PDF 布局完全一致
  3. 封面全面美化：大图分割背景 + 斜切色块 + 精致排版
  4. 性能优化：
     - 流水线处理：EXIF-only 快速扫描 → 聚类 → 按需加载图像
     - 缩略图磁盘缓存（.album_cache/）
     - ThreadPoolExecutor 多线程评分
     - 只处理候选照片（每事件 N 张），大幅减少 I/O
  5. 电影节奏排版（Opening / Narrative / Highlight / Closing）
  6. 故事模板生成（结构化数据驱动，不依赖 AI）

Dependencies:
    pip install pillow reportlab python-pptx exifread
    pip install pillow-heif   # optional HEIC
    exifread is strongly recommended for reliable GPS extraction.
    Without it, GPS falls back to Pillow which may miss GPS on some files.
"""

import os, math, random, threading, hashlib, json, time, sys
from concurrent.futures import ThreadPoolExecutor, as_completed
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from datetime import datetime
from collections import defaultdict
from pathlib import Path
from io import BytesIO

from PIL import Image, ExifTags, ImageOps, ImageFilter, ImageEnhance, ImageStat, ImageDraw
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.units import mm
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont

from everalbum.services.config import AlbumBuildRequest
from everalbum.services.narrative_engine import NarrativeEngine, is_non_location_label
from everalbum.services.portrait_assets import PortraitAsset, PortraitAssetStore
from everalbum.services.workspace import get_default_portrait_library_path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
    HEIC_SUPPORTED = True
except ImportError:
    HEIC_SUPPORTED = False

pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
FONT_REGULAR = "STSong-Light"

SUPPORTED_EXTS = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp"}
if HEIC_SUPPORTED:
    SUPPORTED_EXTS |= {".heic", ".heif"}

# ── A4 dimensions ────────────────────────────────────────────────────────────
PAGE_W, PAGE_H = A4          # 595.28 × 841.89 pt
MARGIN   = 18 * mm
INNER_W  = PAGE_W - 2 * MARGIN
INNER_H  = PAGE_H - 2 * MARGIN

# ── A4 in EMU for PPTX (1 mm = 36000 EMU) ───────────────────────────────────
A4_EMU_W = 7560000   # 210 mm
A4_EMU_H = 10692000  # 297 mm

MONTH_ZH = ["一月","二月","三月","四月","五月","六月",
            "七月","八月","九月","十月","十一月","十二月"]
MONTH_EN = ["January","February","March","April","May","June",
            "July","August","September","October","November","December"]

# ══ SEASONAL PALETTES ════════════════════════════════════════════════════════

def get_season(month):
    if month in (3,4,5):   return "spring"
    if month in (6,7,8):   return "summer"
    if month in (9,10,11): return "autumn"
    return "winter"

SEASON_PALETTES = {
    "spring": {"bg":"#FEF6F0","accent":"#C94B7B","text":"#3B1F2B","sub":"#A06070",
               "cover_dark":"#5A1530","cover_light":"#F9C8D8","tint":(0.85,0.55,0.68)},
    "summer": {"bg":"#F0F8FF","accent":"#0070A8","text":"#012A4A","sub":"#3A7ACC",
               "cover_dark":"#012A4A","cover_light":"#B8DEFF","tint":(0.40,0.65,0.88)},
    "autumn": {"bg":"#FDF3E7","accent":"#B84A10","text":"#2C1A0E","sub":"#9E5010",
               "cover_dark":"#3A1800","cover_light":"#FFCF9A","tint":(0.78,0.40,0.12)},
    "winter": {"bg":"#F2F4F8","accent":"#2A3E5C","text":"#1A232E","sub":"#506080",
               "cover_dark":"#0D1B2A","cover_light":"#C8D8F0","tint":(0.55,0.68,0.85)},
}

MANUAL_PALETTES = {
    "四季自动":    None,
    "暖调 Warm":   {"bg":"#FDF6EC","accent":"#C0392B","text":"#3D2B1F","sub":"#8D6E63",
                   "cover_dark":"#3D2B1F","cover_light":"#F9E0C8","tint":(0.75,0.32,0.18)},
    "冷调 Cool":   {"bg":"#EEF4FB","accent":"#1A6FA0","text":"#1B2A3A","sub":"#4A7FA5",
                   "cover_dark":"#0A1828","cover_light":"#B8D8F0","tint":(0.38,0.58,0.80)},
    "自然 Forest": {"bg":"#F2F7F2","accent":"#2E7D32","text":"#1B3A2B","sub":"#558B5E",
                   "cover_dark":"#0D2A14","cover_light":"#B8E0BC","tint":(0.30,0.58,0.35)},
    "暗调 Dark":   {"bg":"#1E1E2E","accent":"#CBA6F7","text":"#CDD6F4","sub":"#A6ADC8",
                   "cover_dark":"#0E0E1E","cover_light":"#6E6E9E","tint":(0.18,0.16,0.32)},
    "胶片 Film":   {"bg":"#F5F0E8","accent":"#8B6914","text":"#2C2416","sub":"#9E8B70",
                   "cover_dark":"#1E1200","cover_light":"#F0D898","tint":(0.58,0.46,0.20)},
}

def resolve_palette(palette_name, month=None):
    if palette_name == "四季自动":
        return SEASON_PALETTES[get_season(month) if month else "autumn"]
    p = MANUAL_PALETTES.get(palette_name)
    return p if p else SEASON_PALETTES["autumn"]

# ══ UTILITIES ════════════════════════════════════════════════════════════════

def hex_to_rgb(h):
    h=h.lstrip("#"); return tuple(int(h[i:i+2],16)/255 for i in (0,2,4))

def hex_color(h):
    r,g,b=hex_to_rgb(h); return colors.Color(r,g,b)

def hex_rgb(h):
    h=h.lstrip("#"); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)) if HAS_PPTX else None

def dms_to_decimal(dms,ref):
    d,m,s=[x[0]/x[1] for x in dms]; dd=d+m/60+s/3600
    return -dd if ref in ("S","W") else dd

def haversine_km(lat1,lon1,lat2,lon2):
    R=6371; dlat=math.radians(lat2-lat1); dlon=math.radians(lon2-lon1)
    a=math.sin(dlat/2)**2+math.cos(math.radians(lat1))*math.cos(math.radians(lat2))*math.sin(dlon/2)**2
    return R*2*math.asin(math.sqrt(a))

# GPS IFD tag ID
_GPS_IFD = 0x8825

def get_exif(img):
    """Legacy helper — kept for datetime extraction only."""
    try:
        raw = img._getexif()
        return {ExifTags.TAGS.get(k,k): v for k,v in raw.items()} if raw else {}
    except:
        return {}

def extract_datetime(exif):
    for tag in ("DateTimeOriginal","DateTime","DateTimeDigitized"):
        val = exif.get(tag)
        if val:
            try: return datetime.strptime(val, "%Y:%m:%d %H:%M:%S")
            except: pass
    return None

def _dms_from_ifd(val):
    """Convert IFDRational or tuple-pair DMS to float degrees."""
    result = []
    for v in val[:3]:
        if hasattr(v, 'numerator'):        # IFDRational (Pillow 9+)
            result.append(float(v))
        elif isinstance(v, tuple) and len(v) == 2:  # (num, den) pair
            result.append(v[0] / v[1] if v[1] else 0.0)
        else:
            result.append(float(v))
    d, m, s = result
    return d + m/60 + s/3600

def extract_gps(path_or_img):
    """
    Robust GPS extraction using three methods in priority order:
    1. exifread  — same library as user's verified test.py (most reliable)
    2. Modern PIL getexif().get_ifd(0x8825)  — Pillow 9+ clean API
    3. Legacy PIL _getexif() + GPSInfo dict  — old fallback
    Returns (lat, lon) in decimal degrees, or (None, None).
    """
    # Determine file path
    if isinstance(path_or_img, str):
        path = path_or_img
    elif hasattr(path_or_img, 'filename'):
        path = path_or_img.filename
    else:
        path = None

    # ── Method 1: exifread (mirrors test.py exactly) ─────────────────
    if path:
        try:
            import exifread
            with open(path, 'rb') as f:
                tags = exifread.process_file(
                    f, stop_tag='GPS GPSLongitude', details=False)
            if "GPS GPSLatitude" in tags and "GPS GPSLongitude" in tags:
                def conv(tag_val):
                    d = float(tag_val.values[0].num) / float(tag_val.values[0].den)
                    m = float(tag_val.values[1].num) / float(tag_val.values[1].den)
                    s = float(tag_val.values[2].num) / float(tag_val.values[2].den)
                    return d + m/60 + s/3600
                lat = conv(tags["GPS GPSLatitude"])
                lon = conv(tags["GPS GPSLongitude"])
                if str(tags.get("GPS GPSLatitudeRef", "N")).strip() != "N":
                    lat = -lat
                if str(tags.get("GPS GPSLongitudeRef", "E")).strip() != "E":
                    lon = -lon
                return lat, lon
        except ImportError:
            pass  # exifread not installed, try next method
        except Exception:
            pass

    # ── Method 2: Modern PIL getexif().get_ifd() ─────────────────────
    if path:
        try:
            with Image.open(path) as img:
                exif_obj = img.getexif()
                if exif_obj:
                    gps_ifd = exif_obj.get_ifd(_GPS_IFD)
                    if gps_ifd:
                        gps = {ExifTags.GPSTAGS.get(k, k): v
                               for k, v in gps_ifd.items()}
                        if "GPSLatitude" in gps and "GPSLongitude" in gps:
                            lat = _dms_from_ifd(gps["GPSLatitude"])
                            lon = _dms_from_ifd(gps["GPSLongitude"])
                            if str(gps.get("GPSLatitudeRef","N")).strip() not in ("N",""):
                                lat = -lat
                            if str(gps.get("GPSLongitudeRef","E")).strip() not in ("E",""):
                                lon = -lon
                            return lat, lon
        except Exception:
            pass

    # ── Method 3: Legacy PIL _getexif() ──────────────────────────────
    try:
        img_obj = path_or_img if not isinstance(path_or_img, str) else None
        if path and img_obj is None:
            img_obj = Image.open(path)
        if img_obj is not None:
            raw = img_obj._getexif()
            if raw:
                exif_flat = {ExifTags.TAGS.get(k,k): v for k,v in raw.items()}
                gps_raw = exif_flat.get("GPSInfo")
                if gps_raw and isinstance(gps_raw, dict):
                    gps = {ExifTags.GPSTAGS.get(k,k): v
                           for k, v in gps_raw.items()}
                    if "GPSLatitude" in gps and "GPSLongitude" in gps:
                        lat = _dms_from_ifd(gps["GPSLatitude"])
                        lon = _dms_from_ifd(gps["GPSLongitude"])
                        if str(gps.get("GPSLatitudeRef","N")).strip() not in ("N",""):
                            lat = -lat
                        if str(gps.get("GPSLongitudeRef","E")).strip() not in ("E",""):
                            lon = -lon
                        return lat, lon
    except Exception:
        pass

    return None, None

def load_photo(path):
    img=Image.open(path); return ImageOps.exif_transpose(img).convert("RGB")

def resize_fit(img,w,h):
    r=img.width/img.height; br=w/h
    return img.resize((w,max(1,int(w/r))),Image.LANCZOS) if r>br \
           else img.resize((max(1,int(h*r)),h),Image.LANCZOS)

def resize_crop(img,w,h):
    r=img.width/img.height; br=w/h
    nw,nh=(max(w,int(h*r)),h) if r>br else (w,max(h,int(w/r)))
    img=img.resize((nw,nh),Image.LANCZOS)
    return img.crop(((img.width-w)//2,(img.height-h)//2,(img.width-w)//2+w,(img.height-h)//2+h))

def img_to_reader(img,quality=88):
    buf=BytesIO()
    if "A" in img.getbands():
        img.save(buf,format="PNG")
    else:
        img.save(buf,format="JPEG",quality=quality,optimize=True)
    buf.seek(0)
    return ImageReader(buf)

def img_to_bytesio(img,quality=85,max_px=1200):
    thumb=img.copy(); thumb.thumbnail((max_px,max_px),Image.LANCZOS)
    buf=BytesIO()
    if "A" in thumb.getbands():
        thumb.save(buf,format="PNG")
    else:
        thumb.save(buf,format="JPEG",quality=quality)
    buf.seek(0)
    return buf

def shape_polygon_points(shape, width, height):
    if shape == "hexagon":
        return [
            (width * 0.25, 0),
            (width * 0.75, 0),
            (width, height * 0.5),
            (width * 0.75, height),
            (width * 0.25, height),
            (0, height * 0.5),
        ]
    return []

def mask_shape_image(img, w, h, shape="circle", crop=True):
    fitted = resize_crop(img, w, h) if crop else resize_fit(img, w, h)
    rgba = fitted.convert("RGBA")
    mask = Image.new("L", (w, h), 0)
    draw = ImageDraw.Draw(mask)
    if shape == "circle":
        draw.ellipse((0, 0, w - 1, h - 1), fill=255)
    elif shape == "hexagon":
        draw.polygon(shape_polygon_points("hexagon", w - 1, h - 1), fill=255)
    else:
        return rgba
    rgba.putalpha(mask)
    return rgba

def file_hash(path):
    h=hashlib.md5(); h.update(str(os.path.getmtime(path)).encode()); h.update(path.encode()); return h.hexdigest()[:12]

# ══ GEOCODE CACHE ════════════════════════════════════════════════════════════
_geocache: dict = {}          # in-memory: "lat,lon" → place_name
_geocache_path: str = ""      # path to JSON file

def init_geocache(cache_dir: str):
    global _geocache, _geocache_path
    _geocache_path = str(Path(cache_dir) / "geocache.json")
    try:
        with open(_geocache_path, encoding="utf-8") as f:
            _geocache = json.load(f)
    except Exception:
        _geocache = {}

def _save_geocache():
    try:
        with open(_geocache_path, "w", encoding="utf-8") as f:
            json.dump(_geocache, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def _clean_geo_part(s: str) -> str:
    """Take first semicolon-segment and strip whitespace."""
    return s.split(";")[0].strip() if s else ""

def reverse_geocode(lat: float, lon: float) -> str:
    """
    Reverse-geocode lat/lon -> English place name (City, State, Country).
    Uses Nominatim. Results are disk-cached.
    Only called once per event representative photo.
    Returns empty string on failure.
    """
    key = f"{lat:.4f},{lon:.4f}"
    if key in _geocache:
        return _geocache[key]
    try:
        import urllib.request, json as _json
        # English only — avoids "City;Ciudad;Cidade" multilingual artefacts
        url = (f"https://nominatim.openstreetmap.org/reverse?"
               f"lat={lat:.5f}&lon={lon:.5f}&format=json&zoom=10"
               f"&accept-language=en&namedetails=0")
        req = urllib.request.Request(url, headers={"User-Agent": "SmartPhotoAlbum/5.0"})
        with urllib.request.urlopen(req, timeout=6) as r:
            data = _json.loads(r.read())
        addr = data.get("address", {})
        city    = _clean_geo_part(addr.get("city") or addr.get("town")
                                   or addr.get("village") or addr.get("county") or "")
        state   = _clean_geo_part(addr.get("state") or "")
        country = _clean_geo_part(addr.get("country") or "")
        parts = [p for p in [city, state, country] if p]
        place = ", ".join(parts)
        _geocache[key] = place
        _save_geocache()
        return place
    except Exception:
        _geocache[key] = ""
        return ""

def gps_str(lat: float, lon: float) -> str:
    """Format GPS as e.g. 40.7128N 74.0060W"""
    return f"{abs(lat):.4f}{'N' if lat>=0 else 'S'}  {abs(lon):.4f}{'E' if lon>=0 else 'W'}"

def _clean_place_for_display(place: str) -> str:
    """
    Clean a geocoded place name for display.
    Strips anything after the first semicolon (multilingual Nominatim artefact).
    Returns 'Unknown' if empty.
    """
    if not place:
        return "Unknown"
    place = place.split(";")[0].strip()
    place = " ".join(place.split())
    return place or "Unknown"


def default_log(msg):
    text = msg if isinstance(msg, str) else str(msg)
    stream = getattr(sys, "stdout", None)
    encoding = getattr(stream, "encoding", None) or "utf-8"
    try:
        print(text)
    except UnicodeEncodeError:
        safe_text = text.encode(encoding, errors="ignore").decode(encoding, errors="ignore")
        print(safe_text)

# ══ THUMBNAIL CACHE ══════════════════════════════════════════════════════════

class ThumbCache:
    """Disk-based thumbnail cache under .album_cache/ next to first source folder."""
    def __init__(self, cache_dir, thumb_size=1024):
        self.cache_dir = Path(cache_dir)
        self.thumb_size = thumb_size
        self.cache_dir.mkdir(exist_ok=True)

    def _path(self, src_path):
        key = file_hash(src_path)
        return self.cache_dir / f"{key}.jpg"

    def get(self, src_path) -> Image.Image | None:
        cp = self._path(src_path)
        if cp.exists():
            try: return Image.open(cp).convert("RGB")
            except: pass
        return None

    def put(self, src_path, img: Image.Image) -> Image.Image:
        thumb = img.copy()
        thumb.thumbnail((self.thumb_size, self.thumb_size), Image.LANCZOS)
        cp = self._path(src_path)
        try: thumb.save(cp, format="JPEG", quality=85)
        except: pass
        return thumb

    def get_or_make(self, src_path) -> Image.Image:
        cached = self.get(src_path)
        if cached: return cached
        img = load_photo(src_path)
        return self.put(src_path, img)

# ══ MODULE 1: PhotoInfo (EXIF-only fast init) ════════════════════════════════

class PhotoInfo:
    """Lightweight container — image pixels NOT loaded during __init__."""
    __slots__ = ('path','name','dt','lat','lon','location_name','score')

    def __init__(self, path):
        self.path = path
        self.name = Path(path).name
        self.dt: datetime | None = None
        self.lat = self.lon = None
        self.location_name = ""
        self.score: float = 0.3
        self._load_exif_only()

    def _load_exif_only(self):
        """Read EXIF without decoding image pixels — fast even for 5000 files."""
        try:
            with Image.open(self.path) as img:
                exif = get_exif(img)
            self.dt = extract_datetime(exif)
        except: pass
        # GPS: pass file path so extract_gps can use exifread (most reliable)
        try:
            self.lat, self.lon = extract_gps(self.path)
        except: pass
        if not self.dt:
            self.dt = datetime.fromtimestamp(os.path.getmtime(self.path))

    def resolve_location(self):
        """Lazy reverse geocode — call after geocache is initialised."""
        if self.lat is not None and not self.location_name:
            self.location_name = reverse_geocode(self.lat, self.lon)

    @property
    def date_str(self): return self.dt.strftime("%Y年%m月%d日") if self.dt else "未知日期"
    @property
    def time_str(self): return self.dt.strftime("%H:%M") if self.dt else ""
    @property
    def raw_gps_str(self):
        if self.lat is None: return ""
        return gps_str(self.lat, self.lon)
    @property
    def season(self): return get_season(self.dt.month) if self.dt else "autumn"

# ══ MODULE 2: PhotoScorer (threaded, cache-aware) ════════════════════════════

class PhotoScorer:
    """Score photos for sharpness + exposure. Uses thumbnail cache."""

    @staticmethod
    def _sharpness(img):
        gray=img.convert("L").resize((256,256))
        lap=gray.filter(ImageFilter.Kernel(size=3,kernel=[-1,-1,-1,-1,8,-1,-1,-1,-1],scale=1,offset=128))
        return ImageStat.Stat(lap).var[0]

    @staticmethod
    def _exposure(img):
        hist=img.convert("L").histogram(); total=sum(hist)
        if total==0: return 0.5
        dark=sum(hist[:30])/total; bright=sum(hist[225:])/total
        return (1-dark-bright)*0.7+(1-abs(dark-bright))*0.3

    @classmethod
    def score_one(cls, photo: PhotoInfo, cache: ThumbCache):
        try:
            thumb = cache.get_or_make(photo.path)
            sharp = min(cls._sharpness(thumb)/500.0, 1.0)
            photo.score = sharp*0.6 + cls._exposure(thumb)*0.4
        except: photo.score = 0.3

    @classmethod
    def score_batch(cls, photos: list, cache: ThumbCache,
                    max_workers=4, progress_cb=None):
        """Score all photos in parallel using thread pool."""
        total = len(photos); done = 0
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            futures = {ex.submit(cls.score_one, ph, cache): ph for ph in photos}
            for fut in as_completed(futures):
                done += 1
                if progress_cb and done % 20 == 0:
                    progress_cb(int(done/total*100))
        if progress_cb: progress_cb(100)

# ══ MODULE 3: EventClusterer ═════════════════════════════════════════════════

class EventClusterer:
    MIN_EVENT_PHOTOS = 5
    MIN_MONTH_DIVIDER_PHOTOS = 10

    def __init__(self, time_gap_hours=6.0, geo_radius_km=30.0):
        self.time_gap_hours=time_gap_hours; self.geo_radius_km=geo_radius_km

    def cluster(self, photos):
        photos=sorted(photos,key=lambda p:p.dt or datetime.min)
        groups,current=[],[]
        for ph in photos:
            if not current: current.append(ph); continue
            prev=current[-1]
            dt_diff=(abs((ph.dt-prev.dt).total_seconds())/3600 if (ph.dt and prev.dt) else 999)
            geo_diff=(haversine_km(prev.lat,prev.lon,ph.lat,ph.lon) if (prev.lat and ph.lat) else 0)
            if dt_diff>=self.time_gap_hours or geo_diff>=self.geo_radius_km:
                groups.append(current); current=[ph]
            else: current.append(ph)
        if current: groups.append(current)
        return groups

    @staticmethod
    def month_key(group):
        dates = sorted([p.dt for p in group if p.dt])
        return (dates[0].year, dates[0].month) if dates else (0, 0)

    @classmethod
    def merge_small_events(cls, groups, min_photos=None):
        min_photos = min_photos or cls.MIN_EVENT_PHOTOS
        if min_photos <= 1:
            return [sorted(list(group), key=lambda p: p.dt or datetime.min) for group in groups]

        merged = []
        for group in groups:
            current = sorted(list(group), key=lambda p: p.dt or datetime.min)
            current_month = cls.month_key(current)
            if not merged:
                merged.append(current)
                continue
            prev_month = cls.month_key(merged[-1])
            if len(current) < min_photos and current_month == prev_month:
                merged[-1].extend(current)
                merged[-1].sort(key=lambda p: p.dt or datetime.min)
            else:
                merged.append(current)

        result = []
        idx = 0
        while idx < len(merged):
            current = merged[idx]
            current_month = cls.month_key(current)
            if (
                len(current) < min_photos
                and idx + 1 < len(merged)
                and cls.month_key(merged[idx + 1]) == current_month
            ):
                merged[idx + 1] = sorted(current + merged[idx + 1], key=lambda p: p.dt or datetime.min)
            else:
                result.append(current)
            idx += 1
        return result

    @staticmethod
    def month_photo_count(evts):
        return sum(len(grp) for grp, _loc in evts)

    @classmethod
    def should_create_month_page(cls, evts, min_photos=None):
        min_photos = min_photos or cls.MIN_MONTH_DIVIDER_PHOTOS
        return cls.month_photo_count(evts) >= min_photos

    @staticmethod
    def event_label(group):
        """
        Returns clean English place name for the event.
          1. Cleaned location_name (from geocoding stage)
          2. Raw GPS string        (fallback)
          3. Date string           (last resort)
        """
        locs = [p.location_name for p in group if p.location_name]
        if locs:
            return _clean_place_for_display(locs[0])
        gps_cands = [p for p in group if p.lat is not None]
        if gps_cands:
            rep = gps_cands[len(gps_cands)//2]
            return gps_str(rep.lat, rep.lon)
        dates = sorted([p.dt for p in group if p.dt])
        return dates[0].strftime("%Y年%m月%d日") if dates else "Unknown"

    @staticmethod
    def group_by_month(groups, location_names):
        monthly=defaultdict(list)
        for grp,loc in zip(groups,location_names):
            dates=sorted([p.dt for p in grp if p.dt])
            key=(dates[0].year,dates[0].month) if dates else (0,0)
            monthly[key].append((grp,loc))
        return [(y,m,monthly[(y,m)]) for (y,m) in sorted(monthly.keys())]

# ══ MODULE 3b: Pacing Engine (cinematic rhythm) ══════════════════════════════

class PacingEngine:
    """
    Assigns each photo a role: opening / narrative / highlight / closing.
    Drives layout selection for a cinematic feel.
    """
    @staticmethod
    def assign_roles(photos: list) -> list:
        """Returns list of (photo, role) tuples."""
        n = len(photos)
        if n == 0: return []
        if n == 1: return [(photos[0], "highlight")]
        scored = sorted(enumerate(photos), key=lambda x: x[1].score, reverse=True)
        highlight_idx = scored[0][0]
        result = []
        for i, ph in enumerate(photos):
            if i == 0:              role = "opening"
            elif i == n-1:          role = "closing"
            elif i == highlight_idx:role = "highlight"
            else:                   role = "narrative"
            result.append((ph, role))
        return result

    # Track layouts used in the current event to avoid repetition
    _used_layouts: list = []
    _strip_used: bool = False

    @classmethod
    def reset_event(cls):
        """Call at the start of each event to reset layout history."""
        cls._used_layouts = []
        cls._strip_used = False

    @classmethod
    def layout_for_role(cls, role: str, n_available: int) -> str:
        """
        Choose layout name based on cinematic role and available photos.
        Guarantees variety: polaroid/scattered/staircase always appear
        when enough photos are available. Strip is limited to once per event.
        """
        if role == "opening":
            # First page: big establishing shot (occasionally diptych)
            if n_available >= 2 and random.random() < 0.25:
                return "wide_banner"
            return "single"

        if role == "highlight":
            # Best photo: full page or magazine treatment
            if n_available >= 3 and random.random() < 0.30:
                return "magazine"
            return "single"

        if role == "closing":
            if n_available >= 3:
                return random.choice(["magazine","triptych","brick3"])
            if n_available >= 2:
                return random.choice(["diptych","portrait_pair","wide_banner"])
            return "single"

        # ── narrative: rich variety, no repetition ─────────────────────
        # All layouts grouped by minimum photo requirement
        pools = {
            1: ["single"],
            2: ["diptych", "wide_banner", "portrait_pair"],
            3: ["triptych", "brick3", "staircase", "spotlight",
                "asymmetric3", "magazine"],
            4: ["quad", "l_shape", "polaroid", "scattered", "panorama"],
            5: ["mosaic5", "hero_thumbs", "frame5", "cross_center", "circle_cluster"],
            6: ["tiled3x2", "collage6"],  # strip removed here — added below
            7: ["hexagon_hive"],
            9: ["nine_grid"],
        }
        # Build candidate list from all pools where req <= n_available
        candidates = []
        for req, names in sorted(pools.items()):
            if req <= n_available:
                candidates.extend(names)

        # Strip: allow only once per event, and only when ≥6 photos
        if n_available >= 6 and not cls._strip_used:
            candidates.append("strip")

        # Remove layouts already used this event (avoid repetition)
        fresh = [c for c in candidates if c not in cls._used_layouts]
        if not fresh:
            # All layouts used — reset history and pick any
            cls._used_layouts = []
            fresh = candidates

        # Weighted selection: polaroid/scattered/staircase get extra weight
        # so they actually appear in practice
        PRIORITY = {"polaroid":3, "scattered":3, "staircase":2,
                    "frame5":2, "cross_center":2, "collage6":2, "tiled3x2":2,
                    "circle_cluster":2, "hexagon_hive":2, "nine_grid":2}
        weighted = []
        for name in fresh:
            weighted.extend([name] * PRIORITY.get(name, 1))

        chosen = random.choice(weighted)
        cls._used_layouts.append(chosen)
        if chosen == "strip":
            cls._strip_used = True
        return chosen

# ══ SCENE ANALYSER (lightweight, PIL-only) ══════════════════════════════════

class SceneAnalyzer:
    """
    Classifies photos into scene types using colour + brightness analysis.
    No CV2 or PyTorch needed — pure PIL.
    Covers ~80% of common travel photo scenarios.
    """

    @staticmethod
    def _hsv_ratios(img: Image.Image) -> dict:
        """Sample a thumbnail and return colour region ratios."""
        thumb = img.copy(); thumb.thumbnail((128, 128), Image.LANCZOS)
        rgb = thumb.convert("RGB")
        pixels = list(rgb.getdata())
        n = len(pixels)
        sky_blue = water_blue = green = red_warm = dark = bright = 0
        for r, g, b in pixels:
            maxc = max(r, g, b); minc = min(r, g, b)
            lum = (r * 299 + g * 587 + b * 114) / 1000
            if lum < 55:   dark += 1
            if lum > 200:  bright += 1
            # Sky blue: high blue, mid-high brightness, low saturation variation
            if b > 120 and b > r + 20 and b > g - 10 and lum > 120:
                sky_blue += 1
            # Water blue: similar but cooler, can be darker
            if b > 100 and b > r + 15 and g > r and lum > 60:
                water_blue += 1
            # Green vegetation
            if g > 80 and g > r * 1.1 and g > b * 0.9 and lum < 200:
                green += 1
            # Warm reds/oranges (sunset, food, indoor warm light)
            if r > 160 and r > g * 1.2 and r > b * 1.3:
                red_warm += 1
        return {
            "sky_blue":   sky_blue / n,
            "water_blue": water_blue / n,
            "green":      green / n,
            "red_warm":   red_warm / n,
            "dark":       dark / n,
            "bright":     bright / n,
        }

    @staticmethod
    def _top_half_brightness(img: Image.Image) -> float:
        """Average brightness of the top 40% of image (sky detection helper)."""
        w, h = img.size
        top = img.crop((0, 0, w, int(h * 0.4))).convert("L")
        stat = ImageStat.Stat(top)
        return stat.mean[0]

    @classmethod
    def classify(cls, img: Image.Image) -> str:
        """
        Returns one of: sky | night | sunset | water | nature |
                        urban | people | indoor | general
        """
        try:
            r = cls._hsv_ratios(img)
            top_bright = cls._top_half_brightness(img)
            avg_lum = r["bright"] + (1 - r["dark"])

            # Night: very dark overall
            if r["dark"] > 0.55:
                return "night"
            # Sunset/golden hour: warm + moderate brightness
            if r["red_warm"] > 0.30 and r["dark"] < 0.3:
                return "sunset"
            # Sky: top half bright, high sky-blue ratio
            if r["sky_blue"] > 0.28 and top_bright > 150:
                return "sky"
            # Water: water-blue dominant, lower top brightness
            if r["water_blue"] > 0.35 and r["sky_blue"] < 0.25:
                return "water"
            # Nature / green
            if r["green"] > 0.30:
                return "nature"
            # General bright outdoor
            if r["bright"] > 0.40:
                return "outdoor"
            return "general"
        except Exception:
            return "general"

    # Scene → (template_pool, emoji_hint)
    SCENE_CAPTIONS = {
        "sky": [
            "天空很高，时间很慢。",
            "云在走，人在看。",
            "几道航迹轻轻划过蓝天。",
            "头顶是整片自由。",
        ],
        "night": [
            "城市把光藏进了夜里。",
            "灯火点亮了这个夜晚。",
            "夜色深了，记忆还在。",
            "黑暗中有人在等光。",
        ],
        "sunset": [
            "黄昏把天空烧成了橙色。",
            "这一刻，光线是最好的滤镜。",
            "余晖短暂，记录它的颜色。",
            "太阳快落山的时候，一切都是金的。",
        ],
        "water": [
            "水面有光，也有倒影。",
            "海不说话，但你懂它。",
            "波光粼粼，时间流淌。",
            "水的颜色今天格外好看。",
        ],
        "nature": [
            "山很远，但今天走近了一点。",
            "绿色让人安静。",
            "这片风景不需要滤镜。",
            "自然比想象中更宽阔。",
        ],
        "outdoor": [
            "阳光打在每一张脸上。",
            "户外的空气永远比室内好一点。",
            "光线很好，记录也很好。",
        ],
        "general": [
            "此刻已成回忆。",
            "光影之间，时间悄然流逝。",
            "短暂，却值得记录。",
            "每一张照片都是一个停顿。",
            "有些瞬间，只有镜头才记得住。",
        ],
    }


# ══ STORY TEMPLATE ENGINE ════════════════════════════════════════════════════

class StoryEngine:
    """Compatibility wrapper around the richer narrative engine."""

    @classmethod
    def _is_non_location(cls, label: str) -> bool:
        return is_non_location_label(label)

    @classmethod
    def generate(cls, group: list, location_label: str,
                 cache: "ThumbCache | None" = None) -> str:
        scene = "general"
        if group and cache:
            try:
                best_ph = max(group, key=lambda p: p.score)
                thumb = cache.get_or_make(best_ph.path)
                scene = SceneAnalyzer.classify(thumb)
            except Exception:
                scene = "general"
        return NarrativeEngine.generate_from_group(group, location_label, scene=scene)

    @classmethod
    def plan(cls, group: list, location_label: str,
             cache: "ThumbCache | None" = None):
        scene = "general"
        if group and cache:
            try:
                best_ph = max(group, key=lambda p: p.score)
                thumb = cache.get_or_make(best_ph.path)
                scene = SceneAnalyzer.classify(thumb)
            except Exception:
                scene = "general"
        return NarrativeEngine.build_plan_from_group(group, location_label, scene=scene)

# ══ MODULE 4: PDF AlbumBuilder ═══════════════════════════════════════════════

class AlbumBuilder:

    def __init__(self, output_path, palette_name, cache: ThumbCache,
                 progress_cb=None, log_cb=None,
                 portrait_store: PortraitAssetStore | None = None,
                 portrait_assets: list[PortraitAsset] | None = None):
        self.output_path = output_path
        self.palette_name = palette_name
        self.cache = cache
        self.progress_cb = progress_cb or (lambda v:None)
        self.log_cb = log_cb or default_log
        self.portrait_store = portrait_store
        self.portrait_assets = list(portrait_assets or [])
        self.c = rl_canvas.Canvas(output_path, pagesize=A4)
        self.c.setTitle("Photo Album")
        self._page_count = 0
        self._current_pal = None

    def _set_pal(self, month=None):
        self._current_pal = resolve_palette(self.palette_name, month)

    @property
    def pal(self):
        if self._current_pal is None: self._set_pal()
        return self._current_pal

    # ── Primitives ────────────────────────────────────────────────────────
    def _bg(self):
        self.c.setFillColor(hex_color(self.pal["bg"]))
        self.c.rect(0,0,PAGE_W,PAGE_H,fill=1,stroke=0)

    def _accent_bar(self,y,h=2*mm):
        self.c.setFillColor(hex_color(self.pal["accent"]))
        self.c.rect(MARGIN,y,INNER_W,h,fill=1,stroke=0)

    def _text(self,x,y,text,size=10,bold=False,color=None,align="left"):
        color=color or self.pal["text"]
        self.c.setFillColor(hex_color(color)); self.c.setFont(FONT_REGULAR,size)
        if align=="center":
            self.c.drawCentredString(x,y,text)
            if bold: self.c.drawCentredString(x+0.35,y,text)
        elif align=="right":
            self.c.drawRightString(x,y,text)
            if bold: self.c.drawRightString(x+0.35,y,text)
        else:
            self.c.drawString(x,y,text)
            if bold: self.c.drawString(x+0.35,y,text)

    @staticmethod
    def _wrap_text(text: str, max_w: float, size: float) -> list:
        """
        Split text into lines that fit within max_w at given font size.
        CJK chars ≈ size pt wide, ASCII ≈ size*0.55 pt wide.
        Returns list of line strings.
        """
        def char_w(ch):
            return size if ord(ch) > 127 else size * 0.55
        lines, current, cur_w = [], "", 0.0
        for ch in list(text):
            cw = char_w(ch)
            if cur_w + cw > max_w and current:
                lines.append(current); current = ch; cur_w = cw
            else:
                current += ch; cur_w += cw
        if current:
            lines.append(current)
        return lines

    def _text_fit(self, x, y, text: str, max_w: float, max_h: float,
                  size=10, bold=False, color=None, min_size=6) -> float:
        """
        Draw text fitting within (max_w × max_h).
        Algorithm:
          1. Try wrapping at `size`.
          2. If total height > max_h, shrink font in steps until it fits
             (down to min_size).
          3. Draw lines top-to-bottom, return y after last line.
        """
        if not text: return y
        color = color or self.pal["text"]
        cur_size = size

        while cur_size >= min_size:
            line_h  = cur_size * 1.45
            lines   = self._wrap_text(text, max_w, cur_size)
            total_h = len(lines) * line_h
            if total_h <= max_h or cur_size <= min_size:
                break
            cur_size -= 0.5          # shrink by 0.5pt and retry

        # Clamp
        cur_size = max(cur_size, min_size)
        line_h   = cur_size * 1.45
        lines    = self._wrap_text(text, max_w, cur_size)

        self.c.setFillColor(hex_color(color))
        self.c.setFont(FONT_REGULAR, cur_size)
        for line in lines:
            self.c.drawString(x, y, line)
            if bold: self.c.drawString(x + 0.35, y, line)
            y -= line_h
        return y

    # Keep _text_wrapped as thin alias (used in caption_bar via _text)
    def _text_wrapped(self, x, y, text: str, max_w: float, size=10,
                      bold=False, color=None, line_h=None) -> float:
        """Wrap without height constraint (used for captions)."""
        if not text: return y
        if line_h is None: line_h = size * 1.45
        color = color or self.pal["text"]
        self.c.setFillColor(hex_color(color))
        self.c.setFont(FONT_REGULAR, size)
        for line in self._wrap_text(text, max_w, size):
            self.c.drawString(x, y, line)
            if bold: self.c.drawString(x + 0.35, y, line)
            y -= line_h
        return y

    def _place_img_path(self, path, x, y, w, h, crop=True):
        """Load via cache, then place."""
        try:
            thumb = self.cache.get_or_make(path)
            pw,ph = max(1,int(w*2.5)),max(1,int(h*2.5))
            sized = resize_crop(thumb,pw,ph) if crop else resize_fit(thumb,pw,ph)
            self.c.drawImage(img_to_reader(sized),x,y,width=w,height=h,
                             preserveAspectRatio=not crop,mask="auto")
        except: pass

    def _place_image(self,img,x,y,w,h,crop=True):
        pw,ph=max(1,int(w*2.5)),max(1,int(h*2.5))
        thumb=resize_crop(img,pw,ph) if crop else resize_fit(img,pw,ph)
        self.c.drawImage(img_to_reader(thumb),x,y,width=w,height=h,
                         preserveAspectRatio=not crop,mask="auto")

    def _place_shaped_image(self, img, x, y, w, h, shape="circle", crop=True):
        pw, ph = max(1, int(w * 2.5)), max(1, int(h * 2.5))
        shaped = mask_shape_image(img, pw, ph, shape=shape, crop=crop)
        self.c.drawImage(
            img_to_reader(shaped),
            x,
            y,
            width=w,
            height=h,
            preserveAspectRatio=False,
            mask="auto",
        )

    def _new_page(self, month=None):
        if self._page_count>0: self.c.showPage()
        self._set_pal(month); self._bg(); self._page_count+=1

    def _footer(self, page_num):
        self.c.setFillColor(hex_color(self.pal["sub"]))
        self.c.setFont(FONT_REGULAR,7)
        self.c.drawCentredString(PAGE_W/2, 8*mm, f"- {page_num} -")

    def _caption_bar(self, ph, x, y, w):
        """Photo caption: date | time | raw GPS only. No place names here."""
        parts=[]
        if ph.dt:
            parts.append(ph.date_str)
            if ph.time_str: parts.append(ph.time_str)
        # Raw GPS coords only — place names are for chapter pages only
        raw = ph.raw_gps_str
        if raw:
            parts.append(raw)
        caption=" | ".join(parts)
        self.c.setFillColor(hex_color(self.pal["accent"]))
        self.c.circle(x+5,y+4,2.5,fill=1,stroke=0)
        self._text(x+12,y,caption,size=9,color=self.pal["sub"])

    def _page_story_note(self, note: str, role: str):
        if not note:
            return
        role_label = {
            "opening": "OPENING",
            "narrative": "FLOW",
            "highlight": "HIGHLIGHT",
            "closing": "CLOSING",
        }.get(role, "FLOW")
        box_x = MARGIN
        rect_y = PAGE_H - MARGIN + 2.6 * mm
        box_w = INNER_W
        box_h = 9.8 * mm
        self.c.setFillColor(colors.Color(1, 1, 1, 0.72))
        self.c.roundRect(box_x, rect_y, box_w, box_h, 2 * mm, fill=1, stroke=0)
        self.c.setFillColor(hex_color(self.pal["accent"]))
        self.c.roundRect(box_x + 2 * mm, rect_y + 2.3 * mm, 16 * mm, 3.5 * mm, 1.5 * mm, fill=1, stroke=0)
        self._text(box_x + 4 * mm, rect_y + 3.2 * mm, role_label, size=6, color="#FFFFFF")
        self._text_fit(
            box_x + 21 * mm,
            rect_y + 3.2 * mm,
            note,
            max_w=box_w - 23 * mm,
            max_h=6 * mm,
            size=7.2,
            color=self.pal["text"],
            min_size=6,
        )

    def _pick_portrait_asset(self, group, slot="chapter") -> PortraitAsset | None:
        if not self.portrait_assets:
            return None
        seed_parts = [slot, str(len(group))]
        seed_parts.extend(Path(ph.path).name for ph in group[:3])
        return random.Random("|".join(seed_parts)).choice(self.portrait_assets)

    def _draw_portrait_asset(self, asset, x, y, max_w, max_h):
        if not asset or not self.portrait_store:
            return
        try:
            img = self.portrait_store.open_image(asset)
            scale = min(max_w / max(img.width, 1), max_h / max(img.height, 1))
            if scale <= 0:
                return
            draw_w = max(1.0, img.width * scale)
            draw_h = max(1.0, img.height * scale)
            fitted = img.resize(
                (max(1, int(draw_w * 2.5)), max(1, int(draw_h * 2.5))),
                Image.LANCZOS,
            )
            px = x + max_w - draw_w
            self.c.setStrokeColor(colors.Color(*hex_to_rgb(self.pal["accent"]), 0.18))
            self.c.setLineWidth(0.6)
            self.c.line(x, y + 4, x + max_w, y + 4)
            self.c.drawImage(
                img_to_reader(fitted),
                px,
                y,
                width=draw_w,
                height=draw_h,
                mask="auto",
            )
        except Exception:
            pass

    def _draw_chapter_portrait(self, group, x, y, max_w):
        asset = self._pick_portrait_asset(group, slot="chapter")
        if not asset:
            return
        self._draw_portrait_asset(asset, x, y, max_w * 0.92, PAGE_H * 0.26)

    # ══ COVER (v6: cinematic full-bleed with editorial grid) ══════════════
    def build_cover(self, groups, title, monthly_groups):
        """
        Cinematic cover design:
        - Full-page hero photo as background (best score)
        - Right 28%: vertical strip of 3 thumbnail photos with gap
        - Multi-layer gradient: dark bottom half, coloured top vignette
        - Frosted-glass title block (bottom-left quadrant)
        - Bullet city timeline at very bottom
        - Accent colour left stripe + top corner detail
        """
        all_photos = [p for g in groups for p in g]
        dates = sorted([p.dt for p in all_photos if p.dt])
        first_month = dates[0].month if dates else None
        self._new_page(month=first_month)
        pal = self.pal
        dark = pal.get("cover_dark", "#0D1B2A")
        acc  = pal["accent"]
        tr,tg,tb = pal.get("tint", (0.55,0.65,0.80))

        best = sorted(all_photos, key=lambda ph2: ph2.score, reverse=True)

        # ── LAYER 1: Full-bleed hero photo (whole page background) ───────
        if best:
            try:
                img = load_photo(best[0].path)
                # Slightly darken and desaturate for text legibility
                img = ImageEnhance.Brightness(img).enhance(0.70)
                img = ImageEnhance.Color(img).enhance(0.85)
                self._place_image(img, 0, 0, PAGE_W, PAGE_H, crop=True)
            except: pass

        # ── LAYER 2: Dual gradient overlays ──────────────────────────────
        # Bottom-up dark gradient (covers lower 55% for title legibility)
        for i in range(70):
            frac = i / 70
            alpha = frac ** 1.4 * 0.88
            self.c.setFillColor(colors.Color(*hex_to_rgb(dark), alpha))
            self.c.rect(0, 0, PAGE_W, PAGE_H * frac + PAGE_H/70,
                        fill=1, stroke=0)
        # Top-right corner vignette (accent colour)
        for i in range(30):
            frac = i / 30; alpha = (1-frac) * 0.28
            self.c.setFillColor(colors.Color(tr, tg, tb, alpha))
            self.c.rect(PAGE_W*(0.5+frac*0.1), PAGE_H*(0.6+frac*0.1),
                        PAGE_W*(0.5-frac*0.1), PAGE_H*(0.4-frac*0.1),
                        fill=1, stroke=0)

        # ── LAYER 3: Right photo strip (3 photos, right 26%) ─────────────
        STRIP_W = PAGE_W * 0.26
        STRIP_X = PAGE_W - STRIP_W
        STRIP_GAP = 1.5*mm
        strip_photos = best[1:4]
        if len(strip_photos) < 3:
            strip_photos = (strip_photos * 3)[:3]
        # Strip covers top 55% of page
        STRIP_H_TOTAL = PAGE_H * 0.55
        slot_h = (STRIP_H_TOTAL - STRIP_GAP*2) / 3
        strip_base = PAGE_H - STRIP_H_TOTAL
        for i, ph in enumerate(strip_photos):
            try:
                img = load_photo(ph.path)
                img = ImageEnhance.Brightness(img).enhance(0.80)
                sy = strip_base + i*(slot_h + STRIP_GAP)
                self._place_image(img, STRIP_X, sy, STRIP_W, slot_h)
                # Thin white separator line
                if i > 0:
                    self.c.setStrokeColor(colors.Color(1,1,1,0.6))
                    self.c.setLineWidth(0.8)
                    self.c.line(STRIP_X, sy, PAGE_W, sy)
            except: pass
        # Right strip left edge: accent rule
        self.c.setFillColor(hex_color(acc))
        self.c.rect(STRIP_X, strip_base, 2.5*mm, STRIP_H_TOTAL, fill=1, stroke=0)

        # ── LAYER 4: Left accent stripe (full height) ─────────────────────
        self.c.setFillColor(hex_color(acc))
        self.c.rect(0, 0, 3.5*mm, PAGE_H, fill=1, stroke=0)

        # ── LAYER 5: Frosted-glass title block ────────────────────────────
        # Position: left side, lower 40% of page
        TB_X = 8*mm; TB_W = PAGE_W * 0.68
        TB_Y = MARGIN + 30*mm;  TB_H = PAGE_H * 0.38

        # Frosted glass: slightly lighter than dark with blur effect
        # (achieved via layered semi-transparent rects)
        self.c.setFillColor(colors.Color(*hex_to_rgb(dark), 0.70))
        self.c.roundRect(TB_X, TB_Y, TB_W, TB_H, 3*mm, fill=1, stroke=0)
        # Thin accent top border on the block
        self.c.setFillColor(hex_color(acc))
        self.c.roundRect(TB_X, TB_Y+TB_H-2*mm, TB_W, 2*mm, 1*mm, fill=1, stroke=0)

        # Main title (inside block, upper area)
        title_y = TB_Y + TB_H * 0.68
        self.c.setFillColor(colors.white)
        self.c.setFont(FONT_REGULAR, 34)
        self.c.drawString(TB_X + 6*mm, title_y, title)
        self.c.drawString(TB_X + 6*mm + 0.4, title_y, title)  # faux bold

        # Accent underline (matches title width estimate)
        title_w_est = min(len(title) * 17, TB_W - 12*mm)
        self.c.setStrokeColor(hex_color(acc)); self.c.setLineWidth(1.5)
        self.c.line(TB_X+6*mm, title_y-5, TB_X+6*mm+title_w_est, title_y-5)

        # Date range
        if dates:
            span = f"{dates[0].strftime('%Y.%m.%d')}  —  {dates[-1].strftime('%Y.%m.%d')}"
            self.c.setFont(FONT_REGULAR, 9)
            self.c.setFillColor(colors.Color(1,1,1,0.72))
            self.c.drawString(TB_X+6*mm, TB_Y+TB_H*0.38, span)

        # Photo count + event count pills
        n_total = len(all_photos)
        n_events = sum(len(evts) for _,_,evts in monthly_groups)
        pill_y = TB_Y + TB_H*0.15
        for i,(label_text) in enumerate([f"{n_total} 张", f"{n_events} 段旅程"]):
            px = TB_X + 6*mm + i * 28*mm
            self.c.setFillColor(colors.Color(1,1,1,0.15))
            self.c.roundRect(px, pill_y, 24*mm, 6*mm, 1.5*mm, fill=1, stroke=0)
            self.c.setStrokeColor(colors.Color(1,1,1,0.35)); self.c.setLineWidth(0.5)
            self.c.roundRect(px, pill_y, 24*mm, 6*mm, 1.5*mm, fill=0, stroke=1)
            self.c.setFillColor(colors.Color(1,1,1,0.85)); self.c.setFont(FONT_REGULAR, 7)
            self.c.drawCentredString(px+12*mm, pill_y+1.8*mm, label_text)

        # ── LAYER 6: Bottom city bullet timeline ──────────────────────────
        tl_events = []
        for _y, _m, evts in monthly_groups:
            for grp, loc in evts:
                d = sorted([ph2.dt for ph2 in grp if ph2.dt])
                if d:
                    # City only — first comma-separated part
                    place = loc or ""
                    city = place.split(",")[0].strip() if place else ""
                    if not city or StoryEngine._is_non_location(city):
                        city = d[0].strftime("%m/%d")
                    tl_events.append(city[:8])
                if len(tl_events) >= 7: break
            if len(tl_events) >= 7: break

        if tl_events:
            # Single line of city dots
            dot_y = MARGIN + 12*mm
            dot_x = 8*mm
            self.c.setFont(FONT_REGULAR, 7)
            for i, city in enumerate(tl_events):
                # Bullet dot
                self.c.setFillColor(hex_color(acc))
                self.c.circle(dot_x + 2, dot_y + 3, 2, fill=1, stroke=0)
                # City name
                self.c.setFillColor(colors.Color(1,1,1,0.75))
                self.c.drawString(dot_x + 6, dot_y, city)
                # Advance x
                dot_x += len(city) * 5.5 + 14
                if dot_x > PAGE_W - 20*mm:
                    break

        # ── LAYER 7: Top-left corner detail ───────────────────────────────
        # Small geometric accent: two concentric partial circles
        cx, cy, cr = 0, PAGE_H, 18*mm
        self.c.setStrokeColor(colors.Color(*hex_to_rgb(acc), 0.55))
        self.c.setLineWidth(1)
        self.c.arc(cx-cr, cy-cr, cx+cr, cy+cr, 270, 90)
        self.c.setLineWidth(0.5)
        self.c.setStrokeColor(colors.Color(*hex_to_rgb(acc), 0.30))
        self.c.arc(cx-cr*1.6, cy-cr*1.6, cx+cr*1.6, cy+cr*1.6, 270, 90)

    # ══ TOC ═══════════════════════════════════════════════════════════════
    def build_toc(self, monthly_groups, toc_refs):
        first_month=monthly_groups[0][1] if monthly_groups else None
        self._new_page(month=first_month)
        y=PAGE_H-MARGIN-10
        self._text(MARGIN,y,"目  录",size=22,bold=True); y-=6*mm
        self._accent_bar(y,1.5); y-=12
        for mi,(year,month,evts) in enumerate(monthly_groups):
            show_month = EventClusterer.should_create_month_page(evts)
            if y<MARGIN+30: self._new_page(month=month); y=PAGE_H-MARGIN-20
            if show_month:
                mlabel=f"{year}年  {MONTH_ZH[month-1]}" if month>0 else "未知时间"
                self.c.setFillColor(hex_color(self.pal["accent"]))
                self.c.rect(MARGIN,y-1,3*mm,5*mm,fill=1,stroke=0)
                self._text(MARGIN+5*mm,y,mlabel,size=13,bold=True); y-=14
            for ei,(grp,loc) in enumerate(evts):
                if y<MARGIN+20: self._new_page(month=month); y=PAGE_H-MARGIN-20
                pg=toc_refs.get((mi,ei),"")
                self.c.setStrokeColor(hex_color(self.pal["sub"])); self.c.setLineWidth(0.3)
                self.c.setDash(1,3); self.c.line(MARGIN+8*mm+165,y+4,PAGE_W-MARGIN-22,y+4); self.c.setDash()
                self.c.setFillColor(hex_color(self.pal["sub"])); self.c.circle(MARGIN+6*mm,y+3.5,1.5,fill=1,stroke=0)
                # Show location name only — no date/day mixed in
                toc_loc = loc or "Unknown"
                self._text(MARGIN+8.5*mm,y,toc_loc,size=10)
                if pg: self._text(PAGE_W-MARGIN-5,y,str(pg),size=9,color=self.pal["accent"],align="right")
                y-=13
            y-=6 if show_month else 2

    # ══ MONTH DIVIDER ═════════════════════════════════════════════════════
    def build_month_page(self, year, month):
        self._new_page(month=month)
        self.c.setFillColor(hex_color(self.pal["accent"]))
        self.c.rect(0,0,PAGE_W,PAGE_H,fill=1,stroke=0)
        eng=(MONTH_EN[month-1] if month>0 else "").upper()
        self.c.setFillColor(colors.Color(1,1,1,0.07)); self.c.setFont(FONT_REGULAR,90)
        self.c.drawCentredString(PAGE_W/2,PAGE_H*0.22,eng)
        zh=MONTH_ZH[month-1] if month>0 else "未知"
        self.c.setFillColor(colors.white); self.c.setFont(FONT_REGULAR,52)
        self.c.drawCentredString(PAGE_W/2,PAGE_H*0.56+15,zh)
        self.c.drawCentredString(PAGE_W/2+0.4,PAGE_H*0.56+15,zh)
        self.c.setFont(FONT_REGULAR,14); self.c.setFillColor(colors.Color(1,1,1,0.7))
        self.c.drawCentredString(PAGE_W/2,PAGE_H*0.56-18,str(year))
        self.c.setStrokeColor(colors.Color(1,1,1,0.35)); self.c.setLineWidth(1.5)
        self.c.line(MARGIN*2,MARGIN+5,PAGE_W-MARGIN*2,MARGIN+5)

    # ══ CHAPTER / EVENT PAGE ══════════════════════════════════════════════
    def build_chapter_page(self, idx, group, location_label, story_plan=None):
        """
        Fluid right-column layout — each row consumes actual vertical space
        and passes the cursor down. If a row is too long, font shrinks until
        it fits its height budget. Nothing overlaps.

        Right column row budgets (pt):
          Row 2 (location, 17pt)  : max_h = 40  → 2 lines before shrink
          Row 3 (date,     11pt)  : max_h = 22  → 1 line, shrinks if needed
          Row 4 (GPS,       9pt)  : max_h = 16  → 1 line
              ↓ gap 6pt
          Row 6 (story,     8pt)  : max_h = 36  → up to 4 lines
        Bottom guard: MARGIN + 8pt
        """
        dates=sorted([p.dt for p in group if p.dt])
        month=dates[0].month if dates else None
        self._new_page(month=month)
        photos=group[:3]; col_w=PAGE_W*0.45
        if photos:
            sh=PAGE_H/len(photos)
            for i,ph in enumerate(photos):
                try:
                    thumb=self.cache.get_or_make(ph.path)
                    img=ImageEnhance.Brightness(thumb).enhance(0.72)
                    self.c.drawImage(img_to_reader(
                        resize_crop(img,int(col_w*2.5),int(sh*2.5))),
                        0, PAGE_H-sh*(i+1), col_w, sh)
                except: pass

        rx      = col_w + MARGIN
        right_w = PAGE_W - rx - MARGIN   # usable text width
        BOTTOM  = MARGIN + 8             # hard page bottom boundary

        # ── Anchor: chapter number sits at fixed position ──────────────
        ry = PAGE_H / 2 + 40

        # Row 1: Chapter number (fixed, large, decorative — does not flow)
        self.c.setFillColor(hex_color(self.pal["accent"]))
        self.c.setFont(FONT_REGULAR, 60)
        self.c.drawString(rx, ry+80, f"{idx+1:02d}")
        self.c.drawString(rx+0.4, ry+80, f"{idx+1:02d}")

        # Accent rule
        self.c.setFillColor(hex_color(self.pal["accent"]))
        self.c.rect(rx, ry+68, right_w, 1.5, fill=1, stroke=0)

        # ── Fluid cursor starts just below the accent rule ──────────────
        cursor = ry + 52   # top of first text row (y in PDF = upward)

        if story_plan and cursor > BOTTOM:
            cursor = self._text_fit(
                rx, cursor, story_plan.chapter_kicker,
                max_w=right_w, max_h=18,
                size=8.5, color=self.pal["accent"])
            cursor -= 4

        clean_loc = _clean_place_for_display(location_label)
        gps_cands = [p for p in group if p.lat is not None]
        has_named_location = bool(clean_loc) and not StoryEngine._is_non_location(clean_loc)
        title_text = clean_loc if has_named_location else ""

        # Row 2: Title — only use named places as the big heading
        if title_text and cursor > BOTTOM:
            cursor = self._text_fit(
                rx, cursor, title_text,
                max_w=right_w, max_h=40,
                size=17, bold=True)
            cursor -= 6

        rows = []
        if dates:
            dstr = (dates[0].strftime("%Y年%m月%d日")
                    if dates[0].date() == dates[-1].date()
                    else f"{dates[0].strftime('%m月%d日')} - {dates[-1].strftime('%m月%d日')}")
            rows.append(("日期", dstr))
        if gps_cands:
            rep = gps_cands[len(gps_cands)//2]
            rows.append(("坐标", gps_str(rep.lat, rep.lon)))
        if story_plan and story_plan.summary_line:
            rows.append(("片段", story_plan.summary_line))

        label_w = 12 * mm
        for label, value in rows:
            if not value or cursor <= BOTTOM:
                continue
            self._text(rx, cursor, label, size=7.2, color=self.pal["accent"])
            cursor = self._text_fit(
                rx + label_w, cursor, value,
                max_w=right_w - label_w, max_h=20,
                size=9, color=self.pal["sub"])
            cursor -= 3
        cursor -= 3

        # Row 6: Scene story — only use real place names in the copy
        story_loc = clean_loc if has_named_location else ""
        story = story_plan.chapter_story if story_plan else StoryEngine.generate(group, story_loc, self.cache)
        if story and cursor > BOTTOM:
            self._text_fit(
                rx, cursor, story,
                max_w=right_w, max_h=48,
                size=8, color=self.pal["sub"])

        self._draw_chapter_portrait(group, rx, MARGIN + 18*mm, right_w)

    # ══ PHOTO LAYOUTS (12 layouts) ════════════════════════════════════════
    IMG_Y0=MARGIN+22*mm; CAP_Y=MARGIN+18*mm

    def _lp(self, ph, hires=True):
        """
        Load photo for layout rendering.
        hires=True (default): load original file for maximum quality.
        hires=False: use 1024px thumbnail (fast, for cover collage/back cover only).
        The cache is used for scoring; for rendering we always want full resolution
        so the PDF output is sharp even at print quality.
        """
        if hires:
            try:
                return load_photo(ph.path)
            except:
                return self.cache.get_or_make(ph.path)  # fallback to cache
        return self.cache.get_or_make(ph.path)

    def _layout_single(self,photos):
        ph=photos[0]
        try: self._place_image(self._lp(ph),MARGIN,self.IMG_Y0,INNER_W,INNER_H-22*mm)
        except: pass
        self._caption_bar(ph,MARGIN,self.CAP_Y,INNER_W)

    def _layout_diptych(self,photos):
        gap=3*mm; w=(INNER_W-gap)/2; h=INNER_H-22*mm
        for i,ph in enumerate(photos[:2]):
            try: self._place_image(self._lp(ph),MARGIN+i*(w+gap),self.IMG_Y0,w,h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_triptych(self,photos):
        gap=3*mm; lw=INNER_W*0.58; rw=INNER_W-lw-gap; h=INNER_H-22*mm; hh=(h-gap)/2
        try: self._place_image(self._lp(photos[0]),MARGIN,self.IMG_Y0,lw,h)
        except: pass
        rx=MARGIN+lw+gap
        for row,idx in enumerate([2,1]):
            if idx<len(photos):
                try: self._place_image(self._lp(photos[idx]),rx,self.IMG_Y0+(hh+gap)*row,rw,hh)
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_quad(self,photos):
        gap=3*mm; w=(INNER_W-gap)/2; h=(INNER_H-22*mm-gap)/2
        for i,ph in enumerate(photos[:4]):
            col,row=[(0,1),(1,1),(0,0),(1,0)][i]
            try: self._place_image(self._lp(ph),MARGIN+col*(w+gap),self.IMG_Y0+row*(h+gap),w,h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_mosaic5(self,photos):
        gap=3*mm; th=(INNER_H-22*mm)*0.52; bh=INNER_H-22*mm-th-gap
        tw=(INNER_W-gap)/2; bw=(INNER_W-2*gap)/3; ty2=self.IMG_Y0+bh+gap
        for i,ph in enumerate(photos[:2]):
            try: self._place_image(self._lp(ph),MARGIN+i*(tw+gap),ty2,tw,th)
            except: pass
        for i,ph in enumerate(photos[2:5]):
            try: self._place_image(self._lp(ph),MARGIN+i*(bw+gap),self.IMG_Y0,bw,bh)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_hero_thumbs(self,photos):
        gap=3*mm; hh=(INNER_H-22*mm)*0.65; th=INNER_H-22*mm-hh-gap; hy=self.IMG_Y0+th+gap
        try: self._place_image(self._lp(photos[0]),MARGIN,hy,INNER_W,hh)
        except: pass
        n=min(len(photos)-1,4)
        if n>0:
            tw=(INNER_W-gap*(n-1))/n
            for i,ph in enumerate(photos[1:1+n]):
                try: self._place_image(self._lp(ph),MARGIN+i*(tw+gap),self.IMG_Y0,tw,th)
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_strip(self,photos):
        n=min(len(photos),6); gap=2*mm; w=(INNER_W-gap*(n-1))/n; h=INNER_H-22*mm
        for i,ph in enumerate(photos[:n]):
            try: self._place_image(self._lp(ph),MARGIN+i*(w+gap),self.IMG_Y0,w,h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_wide_banner(self,photos):
        ph=photos[0]; h=INNER_H-22*mm; bh=h*0.70; sh=h-bh-3*mm; by2=self.IMG_Y0+sh+3*mm
        try: self._place_image(self._lp(ph),MARGIN,by2,INNER_W,bh,crop=False)
        except: pass
        if len(photos)>1:
            try: self._place_image(self._lp(photos[1]),MARGIN,self.IMG_Y0,INNER_W,sh)
            except: pass
        self._caption_bar(ph,MARGIN,self.CAP_Y,INNER_W)

    def _layout_l_shape(self,photos):
        gap=3*mm; h=INNER_H-22*mm; lw=INNER_W*0.45; rw=INNER_W-lw-gap; th=h*0.55; bh=h-th-gap
        try: self._place_image(self._lp(photos[0]),MARGIN,self.IMG_Y0,lw,h)
        except: pass
        rx=MARGIN+lw+gap
        if len(photos)>1:
            try: self._place_image(self._lp(photos[1]),rx,self.IMG_Y0+bh+gap,rw,th)
            except: pass
        hw=(rw-gap)/2
        for j in range(2):
            if len(photos)>2+j:
                try: self._place_image(self._lp(photos[2+j]),rx+j*(hw+gap),self.IMG_Y0,hw,bh)
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_polaroid(self,photos):
        n=min(len(photos),4); configs=[(0.05,0.10,-4),(0.52,0.06,+3),(0.02,0.52,+5),(0.50,0.50,-3)]
        pw=INNER_W*0.42; ph_h=pw*1.18; bord=5*mm; iw=pw-2*bord; ih=ph_h-2*bord-14*mm
        for i in range(n):
            fx,fy,rot=configs[i]; px=MARGIN+fx*INNER_W; py=self.IMG_Y0+fy*(INNER_H-22*mm)
            self.c.saveState()
            self.c.translate(px+pw/2,py+ph_h/2); self.c.rotate(rot); self.c.translate(-pw/2,-ph_h/2)
            self.c.setFillColor(colors.white); self.c.setStrokeColor(colors.Color(0.8,0.8,0.8))
            self.c.setLineWidth(0.5); self.c.rect(0,0,pw,ph_h,fill=1,stroke=1)
            try: self._place_image(self._lp(photos[i]),bord,bord+12*mm,iw,ih)
            except: pass
            self.c.setFillColor(hex_color(self.pal["sub"])); self.c.setFont(FONT_REGULAR,7)
            self.c.drawCentredString(pw/2,5,photos[i].date_str)
            self.c.restoreState()
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_asymmetric3(self,photos):
        gap=3*mm; h=INNER_H-22*mm; sw=INNER_W*0.22; mw=INNER_W-2*sw-2*gap
        try: self._place_image(self._lp(photos[0]),MARGIN+sw+gap,self.IMG_Y0,mw,h)
        except: pass
        for j,ph in enumerate(photos[1:3]):
            try: self._place_image(self._lp(ph),MARGIN+j*(sw+mw+2*gap),self.IMG_Y0,sw,h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_magazine(self,photos):
        gap=3*mm; h=INNER_H-22*mm; hh=h*0.62; bh=h-hh-gap; bw=(INNER_W-gap)/2; hy2=self.IMG_Y0+bh+gap
        try: self._place_image(self._lp(photos[0]),MARGIN,hy2,INNER_W,hh)
        except: pass
        for j,ph in enumerate(photos[1:3]):
            try: self._place_image(self._lp(ph),MARGIN+j*(bw+gap),self.IMG_Y0,bw,bh)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)


    def _layout_brick3(self,photos):
        """1 tall left (60%), 2 stacked right (40%) at golden ratio split."""
        gap=3*mm; h=INNER_H-22*mm; lw=INNER_W*0.60; rw=INNER_W-lw-gap
        top_h=h*0.618; bot_h=h-top_h-gap
        try: self._place_image(self._lp(photos[0]),MARGIN,self.IMG_Y0,lw,h)
        except: pass
        rx=MARGIN+lw+gap
        if len(photos)>1:
            try: self._place_image(self._lp(photos[1]),rx,self.IMG_Y0+bot_h+gap,rw,top_h)
            except: pass
        if len(photos)>2:
            try: self._place_image(self._lp(photos[2]),rx,self.IMG_Y0,rw,bot_h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_panorama(self,photos):
        """1 wide top (55% height) + 3 equal columns bottom."""
        gap=3*mm; h=INNER_H-22*mm; top_h=h*0.55; bot_h=h-top_h-gap
        n=min(len(photos)-1,3); bw=(INNER_W-gap*(n-1))/max(n,1)
        try: self._place_image(self._lp(photos[0]),MARGIN,self.IMG_Y0+bot_h+gap,INNER_W,top_h,crop=False)
        except: pass
        for i,ph in enumerate(photos[1:1+n]):
            try: self._place_image(self._lp(ph),MARGIN+i*(bw+gap),self.IMG_Y0,bw,bot_h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_frame5(self,photos):
        """Centre large + 4 smaller corners."""
        gap=3*mm; h=INNER_H-22*mm
        # Corner size
        cw=(INNER_W-3*gap)/3*0.72; ch=cw*0.75
        # Centre occupies remaining middle area
        mw=INNER_W-2*(cw+gap); mh=h-2*(ch+gap); mx=MARGIN+cw+gap; my=self.IMG_Y0+ch+gap
        try: self._place_image(self._lp(photos[0]),mx,my,mw,mh)
        except: pass
        corners=[(MARGIN,self.IMG_Y0),(MARGIN+cw+gap+mw+gap,self.IMG_Y0),
                 (MARGIN,self.IMG_Y0+ch+gap+mh+gap),(MARGIN+cw+gap+mw+gap,self.IMG_Y0+ch+gap+mh+gap)]
        for i,(cx,cy) in enumerate(corners):
            if i+1<len(photos):
                try: self._place_image(self._lp(photos[i+1]),cx,cy,cw,ch)
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_staircase(self,photos):
        """3 photos in diagonal staircase — each offset right and down."""
        gap=3*mm; h=INNER_H-22*mm
        pw=INNER_W*0.54; ph2=h*0.54
        offsets=[(0,h-ph2),(INNER_W*0.23,h*0.23),(INNER_W*0.46,0)]
        for i,(ox,oy) in enumerate(offsets):
            if i<len(photos):
                try: self._place_image(self._lp(photos[i]),MARGIN+ox,self.IMG_Y0+oy,pw,ph2)
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_collage6(self,photos):
        """6-photo editorial collage: 2 large + 4 small."""
        gap=3*mm; h=INNER_H-22*mm
        # Top row: 2 unequal photos (60/40 split)
        top_h=h*0.50; bot_h=h-top_h-gap
        lw2=INNER_W*0.58; rw2=INNER_W-lw2-gap
        try: self._place_image(self._lp(photos[0]),MARGIN,self.IMG_Y0+bot_h+gap,lw2,top_h)
        except: pass
        if len(photos)>1:
            try: self._place_image(self._lp(photos[1]),MARGIN+lw2+gap,self.IMG_Y0+bot_h+gap,rw2,top_h)
            except: pass
        # Bottom row: 4 equal
        bw=(INNER_W-3*gap)/4
        for i,ph in enumerate(photos[2:6]):
            try: self._place_image(self._lp(ph),MARGIN+i*(bw+gap),self.IMG_Y0,bw,bot_h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_spotlight(self,photos):
        """Centre spotlight: large centre, 2 flanking side panels."""
        gap=3*mm; h=INNER_H-22*mm; side_w=INNER_W*0.18; mid_w=INNER_W-2*(side_w+gap)
        try: self._place_image(self._lp(photos[0]),MARGIN+side_w+gap,self.IMG_Y0,mid_w,h)
        except: pass
        for i,ph in enumerate(photos[1:3]):
            x=MARGIN if i==0 else MARGIN+side_w+gap+mid_w+gap
            try: self._place_image(self._lp(ph),x,self.IMG_Y0,side_w,h)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_tiled3x2(self,photos):
        """Clean 3×2 grid of 6 equal tiles."""
        gap=2.5*mm; h=INNER_H-22*mm
        cols,rows=3,2; tw=(INNER_W-gap*(cols-1))/cols; th=(h-gap*(rows-1))/rows
        for i,ph in enumerate(photos[:6]):
            col=i%cols; row=i//cols
            try: self._place_image(self._lp(ph),MARGIN+col*(tw+gap),self.IMG_Y0+row*(th+gap),tw,th)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_nine_grid(self,photos):
        """3×3 contact-sheet grid for dense story moments."""
        gap=2.2*mm; h=INNER_H-22*mm
        cols=rows=3; tw=(INNER_W-gap*(cols-1))/cols; th=(h-gap*(rows-1))/rows
        for i,ph in enumerate(photos[:9]):
            col=i%cols; row=i//cols
            try: self._place_image(self._lp(ph),MARGIN+col*(tw+gap),self.IMG_Y0+row*(th+gap),tw,th)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_circle_cluster(self,photos):
        """One large circular portrait plus four smaller circular echoes."""
        h=INNER_H-22*mm; gap=4*mm
        big_d=min(INNER_W*0.48,h*0.78)
        remaining=INNER_W-big_d-gap
        small_d=min((remaining-gap)/2,(h-gap)/2)
        big_x=MARGIN+2*mm; big_y=self.IMG_Y0+(h-big_d)/2
        slots=[
            (big_x,big_y,big_d,big_d,0),
            (MARGIN+big_d+gap,self.IMG_Y0+h-small_d,small_d,small_d,1),
            (MARGIN+big_d+gap+small_d+gap,self.IMG_Y0+h-small_d,small_d,small_d,2),
            (MARGIN+big_d+gap,self.IMG_Y0,small_d,small_d,3),
            (MARGIN+big_d+gap+small_d+gap,self.IMG_Y0,small_d,small_d,4),
        ]
        for x2,y2,w2,h2,idx in slots:
            if idx < len(photos):
                try: self._place_shaped_image(self._lp(photos[idx]),x2,y2,w2,h2,shape="circle")
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_hexagon_hive(self,photos):
        """Hexagon hive: one hero hexagon with six supporting cells."""
        h=INNER_H-22*mm
        gap=2.6*mm
        big_w=min(INNER_W*0.33,h*0.34); big_h=big_w*0.88
        small_w=big_w*0.74; small_h=small_w*0.88
        cx=MARGIN+INNER_W/2; cy=self.IMG_Y0+h/2
        x_radius=big_w*0.78 + small_w*0.56 + gap
        y_radius=big_h*0.86 + small_h*0.56 + gap
        diag_x=x_radius*0.88
        diag_y=y_radius*0.52

        def centered(center_x, center_y, w2, h2):
            return center_x-w2/2, center_y-h2/2, w2, h2

        slots=[
            centered(cx, cy, big_w, big_h),
            centered(cx, cy+y_radius, small_w, small_h),
            centered(cx+diag_x, cy+diag_y, small_w, small_h),
            centered(cx+diag_x, cy-diag_y, small_w, small_h),
            centered(cx, cy-y_radius, small_w, small_h),
            centered(cx-diag_x, cy-diag_y, small_w, small_h),
            centered(cx-diag_x, cy+diag_y, small_w, small_h),
        ]
        for idx,(x2,y2,w2,h2) in enumerate(slots):
            if idx < len(photos):
                try: self._place_shaped_image(self._lp(photos[idx]),x2,y2,w2,h2,shape="hexagon")
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_portrait_pair(self,photos):
        """2 tall portrait-oriented photos side by side with generous gap."""
        gap=6*mm; w=(INNER_W-gap)/2; h=INNER_H-22*mm
        for i,ph in enumerate(photos[:2]):
            try: self._place_image(self._lp(ph),MARGIN+i*(w+gap),self.IMG_Y0,w,h,crop=True)
            except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_scattered(self,photos):
        """4 photos scattered at slight angles, overlapping polaroid style."""
        n=min(len(photos),4)
        # configs: (x_frac, y_frac, width_frac, height_frac, rotation_deg)
        configs=[
            (0.00, 0.05, 0.52, 0.62, -3),
            (0.45, 0.00, 0.52, 0.55, +4),
            (0.05, 0.45, 0.45, 0.52, +2),
            (0.50, 0.44, 0.48, 0.54, -5),
        ]
        h=INNER_H-22*mm
        for i in range(n):
            fx,fy,fw,fh,rot=configs[i]
            px=MARGIN+fx*INNER_W; py=self.IMG_Y0+fy*h
            pw=fw*INNER_W; ph2=fh*h
            self.c.saveState()
            self.c.translate(px+pw/2,py+ph2/2); self.c.rotate(rot); self.c.translate(-pw/2,-ph2/2)
            # White border
            self.c.setFillColor(colors.white); self.c.setLineWidth(0.3)
            self.c.setStrokeColor(colors.Color(0.7,0.7,0.7))
            border=3*mm
            self.c.rect(-border,-border,pw+2*border,ph2+2*border,fill=1,stroke=1)
            try: self._place_image(self._lp(photos[i]),0,0,pw,ph2)
            except: pass
            self.c.restoreState()
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    def _layout_cross_center(self,photos):
        """Centre square hero + 4 thin strips on each side."""
        gap=2*mm; h=INNER_H-22*mm
        strip=INNER_W*0.15; centre_w=INNER_W-2*(strip+gap); centre_h=h-2*(strip+gap)
        cx=MARGIN+strip+gap; cy=self.IMG_Y0+strip+gap
        try: self._place_image(self._lp(photos[0]),cx,cy,centre_w,centre_h)
        except: pass
        sides=[
            (MARGIN,self.IMG_Y0,strip,h,1),                      # left
            (cx+centre_w+gap,self.IMG_Y0,strip,h,2),              # right
            (cx,self.IMG_Y0,centre_w,strip,3),                    # top
            (cx,self.IMG_Y0+strip+gap+centre_h+gap,centre_w,strip,4), # bottom
        ]
        for x2,y2,w2,h2,idx in sides:
            if idx<len(photos):
                try: self._place_image(self._lp(photos[idx]),x2,y2,w2,h2)
                except: pass
        self._caption_bar(photos[0],MARGIN,self.CAP_Y,INNER_W)

    LAYOUTS=[
        ("single",1),("diptych",2),("portrait_pair",2),("triptych",3),
        ("brick3",3),("staircase",3),("spotlight",3),("asymmetric3",3),
        ("quad",4),("l_shape",4),("polaroid",4),("scattered",4),
        ("mosaic5",5),("hero_thumbs",5),("frame5",5),("circle_cluster",5),
        ("strip",6),("tiled3x2",6),("collage6",6),("hexagon_hive",7),
        ("nine_grid",9),
        ("wide_banner",2),("magazine",3),("panorama",4),("cross_center",5),
    ]

    def _choose_layout(self,n):
        cands=[(nm,req) for nm,req in self.LAYOUTS if req<=n] or [("single",1)]
        total=sum(r for _,r in cands); r=random.random()*total; cum=0
        for nm,req in cands:
            cum+=req
            if r<=cum: return nm,req
        return cands[-1]

    def _render_layout(self,name,photos,page_note="",role="narrative"):
        {
         "single":self._layout_single,
         "diptych":self._layout_diptych,
         "portrait_pair":self._layout_portrait_pair,
         "triptych":self._layout_triptych,
         "brick3":self._layout_brick3,
         "staircase":self._layout_staircase,
         "spotlight":self._layout_spotlight,
         "asymmetric3":self._layout_asymmetric3,
         "quad":self._layout_quad,
         "l_shape":self._layout_l_shape,
         "polaroid":self._layout_polaroid,
         "scattered":self._layout_scattered,
         "mosaic5":self._layout_mosaic5,
         "hero_thumbs":self._layout_hero_thumbs,
         "frame5":self._layout_frame5,
         "circle_cluster":self._layout_circle_cluster,
         "strip":self._layout_strip,
         "tiled3x2":self._layout_tiled3x2,
         "collage6":self._layout_collage6,
         "hexagon_hive":self._layout_hexagon_hive,
         "nine_grid":self._layout_nine_grid,
         "wide_banner":self._layout_wide_banner,
         "magazine":self._layout_magazine,
         "panorama":self._layout_panorama,
         "cross_center":self._layout_cross_center,
         }.get(name,self._layout_single)(photos)
        self._page_story_note(page_note, role)

    # ══ BACK COVER ════════════════════════════════════════════════════════
    def build_back_cover(self, all_photos, album_title):
        self._new_page(month=None)
        random.seed(42)
        sample=sorted(all_photos,key=lambda p:p.score,reverse=True)[:30]
        if len(sample)<6: sample=sample*8
        cols,rows=5,6; cw,rh=PAGE_W/cols,PAGE_H/rows
        positions=[(c,r) for r in range(rows) for c in range(cols)]
        random.shuffle(positions)
        for (col,row),ph in zip(positions,sample):
            try:
                thumb=self.cache.get_or_make(ph.path)
                img=ImageEnhance.Brightness(thumb).enhance(random.uniform(0.52,0.78))
                self._place_image(img,col*cw,row*rh,cw,rh,crop=True)
            except: pass
        # No colour overlay — show original photos at full colour
        bw,bh=140*mm,28*mm; bx=(PAGE_W-bw)/2; by=(PAGE_H-bh)/2
        # Semi-transparent dark badge only (no tint wash)
        self.c.setFillColor(colors.Color(0,0,0,0.55))
        self.c.roundRect(bx,by,bw,bh,4*mm,fill=1,stroke=0)
        self.c.setFillColor(colors.white); self.c.setFont(FONT_REGULAR,22)
        self.c.drawCentredString(PAGE_W/2,by+bh/2+5,album_title)
        self.c.drawCentredString(PAGE_W/2+0.4,by+bh/2+5,album_title)
        self.c.setStrokeColor(hex_color(self.pal["accent"])); self.c.setLineWidth(1.5)
        self.c.line(PAGE_W/2-35*mm,by+bh/2-3,PAGE_W/2+35*mm,by+bh/2-3)

    # ══ MAIN BUILD (with pacing) ═══════════════════════════════════════════
    def build(self, groups, location_names, album_title="我的相册"):
        NarrativeEngine.reset_album_memory()
        all_photos=[p for g in groups for p in g]; total=len(all_photos); done=0
        monthly_groups=EventClusterer.group_by_month(groups,location_names)
        # Pass 1: page number simulation
        toc_refs={}; sim=2  # cover + toc
        ev_idx=0
        for mi,(year,month,evts) in enumerate(monthly_groups):
            if EventClusterer.should_create_month_page(evts):
                sim+=1  # month divider
            for ei,(grp,loc) in enumerate(evts):
                sim+=1; toc_refs[(mi,ei)]=sim  # chapter page
                PacingEngine.reset_event()
                pacing=PacingEngine.assign_roles(list(grp))
                random.seed(ev_idx*137); i2=0
                while i2<len(pacing):
                    ph,role=pacing[i2]
                    n_avail=len(pacing)-i2
                    layout_name=PacingEngine.layout_for_role(role,n_avail)
                    req=next((r for nm,r in self.LAYOUTS if nm==layout_name),1)
                    sim+=1; i2+=max(req,1)
                ev_idx+=1
        # Pass 2: draw
        self.log_cb("🎨 生成封面..."); self.build_cover(groups,album_title,monthly_groups)
        self.log_cb("📋 生成目录..."); self.build_toc(monthly_groups,toc_refs)
        page_num=3; ev_idx=0
        for mi,(year,month,evts) in enumerate(monthly_groups):
            show_month = EventClusterer.should_create_month_page(evts)
            if show_month:
                self.log_cb(f"📅 {year}年 {MONTH_ZH[month-1] if month>0 else '?'}")
                self.build_month_page(year,month); page_num+=1
            for ei,(grp,loc) in enumerate(evts):
                self.log_cb(f"  📖 {loc or '旅途中'}")
                story_plan = StoryEngine.plan(grp, loc, self.cache)
                self.build_chapter_page(ev_idx,grp,loc,story_plan=story_plan); page_num+=1
                PacingEngine.reset_event()   # reset layout history per event
                pacing=PacingEngine.assign_roles(list(grp))
                random.seed(ev_idx*137); i2=0; page_idx=0
                while i2<len(pacing):
                    chunk_pacing=pacing[i2:]
                    role=chunk_pacing[0][1]; n_avail=len(chunk_pacing)
                    layout_name=PacingEngine.layout_for_role(role,n_avail)
                    req=next((r for nm,r in self.LAYOUTS if nm==layout_name),1)
                    page_photos=[ph for ph,_ in chunk_pacing[:max(req,1)]]
                    page_note = story_plan.note_for_role(role, page_idx)
                    self._new_page(month=month)
                    self._render_layout(layout_name,page_photos,page_note=page_note,role=role)
                    self._footer(page_num); page_num+=1
                    step=max(req,1); i2+=step; done+=step; page_idx+=1
                    self.progress_cb(int(done/max(total,1)*85))
                ev_idx+=1
        self.log_cb("📕 生成封底..."); self.build_back_cover(all_photos,album_title)
        self.c.save(); self.log_cb(f"PDF 已保存: {self.output_path}"); self.progress_cb(100)


# ══ MODULE 5: PptxBuilder — A4 size, mirrors PDF layout exactly ══════════════

class PptxBuilder:
    """
    Issue 2 fix: A4 format (210×297mm), layout mirrors PDF exactly.
    No widescreen 16:9.
    """
    # A4 in EMU
    SW = Emu(A4_EMU_W) if HAS_PPTX else None
    SH = Emu(A4_EMU_H) if HAS_PPTX else None
    # Margin in EMU
    MG = Emu(int(MARGIN / mm * 36000))  # MARGIN mm → EMU

    def __init__(self, output_path, palette_name, cache: ThumbCache, log_cb=None,
                 portrait_store: PortraitAssetStore | None = None,
                 portrait_assets: list[PortraitAsset] | None = None):
        self.output_path = output_path
        self.palette_name = palette_name
        self.cache = cache
        self.log_cb = log_cb or default_log
        self.portrait_store = portrait_store
        self.portrait_assets = list(portrait_assets or [])

    def _pal(self, month=None): return resolve_palette(self.palette_name, month)
    def _rgb(self, h): return hex_rgb(h)

    def _pt(self, mm_val):
        """Convert mm to EMU."""
        return Emu(int(mm_val * 36000))

    def _bg_rect(self, slide, color_hex):
        sh=slide.shapes.add_shape(1,0,0,self.SW,self.SH)
        sh.fill.solid(); sh.fill.fore_color.rgb=self._rgb(color_hex); sh.line.fill.background()

    def _textbox(self, slide, text, l, t, w, h, size=12, bold=False,
                 color="000000", align=None):
        if align is None: align = PP_ALIGN.LEFT
        tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=str(text)
        r.font.size=Pt(size); r.font.bold=bold; r.font.name="Arial"
        r.font.color.rgb=self._rgb(color)
        return tb

    def _add_pic(self, slide, ph, l, t, w, h):
        try:
            thumb=self.cache.get_or_make(ph.path)
            buf=img_to_bytesio(thumb,max_px=1200)
            slide.shapes.add_picture(buf,l,t,w,h)
        except: pass

    def _add_pic_masked(self, slide, ph, l, t, w, h, shape="circle"):
        try:
            thumb = self.cache.get_or_make(ph.path)
            px_w = max(1, int(w / 9525))
            px_h = max(1, int(h / 9525))
            shaped = mask_shape_image(thumb, px_w, px_h, shape=shape, crop=True)
            buf = img_to_bytesio(shaped, max_px=max(px_w, px_h, 1200))
            slide.shapes.add_picture(buf, l, t, w, h)
        except: pass

    def _add_img(self, slide, img, l, t, w, h):
        try:
            buf=img_to_bytesio(img,max_px=1200)
            slide.shapes.add_picture(buf,l,t,w,h)
        except: pass

    def _pick_portrait_asset(self, group, slot="chapter") -> PortraitAsset | None:
        if not self.portrait_assets:
            return None
        seed_parts = [slot, str(len(group))]
        seed_parts.extend(Path(ph.path).name for ph in group[:3])
        return random.Random("|".join(seed_parts)).choice(self.portrait_assets)

    def _add_portrait_asset(self, slide, group, left, bottom, max_w, max_h):
        asset = self._pick_portrait_asset(group, slot="chapter")
        if not asset or not self.portrait_store:
            return
        try:
            img = self.portrait_store.open_image(asset)
            scale = min(max_w / max(img.width, 1), max_h / max(img.height, 1))
            if scale <= 0:
                return
            draw_w = max(1, int(img.width * scale))
            draw_h = max(1, int(img.height * scale))
            fitted = img.resize(
                (max(1, int(draw_w * 1.8)), max(1, int(draw_h * 1.8))),
                Image.LANCZOS,
            )
            buf = img_to_bytesio(fitted, max_px=1600)
            slide.shapes.add_picture(buf, left + max_w - draw_w, bottom - draw_h, draw_w, draw_h)
        except Exception:
            pass

    def _slide_story_note(self, slide, note, role, month):
        if not note:
            return
        pal = self._pal(month)
        role_label = {
            "opening": "OPENING",
            "narrative": "FLOW",
            "highlight": "HIGHLIGHT",
            "closing": "CLOSING",
        }.get(role, "FLOW")
        box_x = self._mm(MARGIN / mm)
        box_y = self._mm(8)
        box_w = self.SW - 2 * box_x
        box_h = self._mm(10.5)
        badge = slide.shapes.add_shape(1, box_x, box_y, box_w, box_h)
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
        badge.fill.transparency = 28
        badge.line.fill.background()
        accent = slide.shapes.add_shape(1, box_x + self._mm(2), box_y + self._mm(2), self._mm(16), self._mm(4))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._rgb(pal["accent"])
        accent.line.fill.background()
        self._textbox(slide, role_label, box_x + self._mm(3.8), box_y + self._mm(1.9), self._mm(12), self._mm(4),
                      size=6, bold=False, color="FFFFFF")
        self._textbox(slide, note, box_x + self._mm(21), box_y + self._mm(1.3), box_w - self._mm(22.5), self._mm(7),
                      size=7, color=pal["text"].lstrip("#"))

    def build(self, groups, location_names, album_title="我的相册"):
        if not HAS_PPTX:
            self.log_cb("python-pptx 未安装，跳过 PPTX。pip install python-pptx"); return
        NarrativeEngine.reset_album_memory()
        prs=Presentation()
        prs.slide_width=self.SW; prs.slide_height=self.SH
        blank=prs.slide_layouts[6]
        all_photos=[p for g in groups for p in g]
        monthly_groups=EventClusterer.group_by_month(groups,location_names)
        dates=sorted([p.dt for p in all_photos if p.dt])

        self.log_cb("PPTX: 封面...")
        self._slide_cover(prs,blank,all_photos,album_title,dates,monthly_groups)
        self.log_cb("PPTX: 目录...")
        self._slide_toc(prs,blank,monthly_groups)

        ev_idx=0
        for mi,(year,month,evts) in enumerate(monthly_groups):
            if EventClusterer.should_create_month_page(evts):
                self._slide_month(prs,blank,year,month)
            for ei,(grp,loc) in enumerate(evts):
                story_plan = StoryEngine.plan(grp, loc, self.cache)
                self._slide_chapter(prs,blank,ev_idx,grp,loc,month,story_plan=story_plan)
                # Photo pages using pacing
                PacingEngine.reset_event()
                pacing=PacingEngine.assign_roles(list(grp))
                i2=0; page_idx=0
                while i2<len(pacing):
                    chunk=pacing[i2:]
                    role=chunk[0][1]; n_avail=len(chunk)
                    layout_name=PacingEngine.layout_for_role(role,n_avail)
                    req=next((r for nm,r in AlbumBuilder.LAYOUTS if nm==layout_name),1)
                    page_photos=[ph for ph,_ in chunk[:max(req,1)]]
                    page_note = story_plan.note_for_role(role, page_idx)
                    self._slide_photos(prs,blank,page_photos,month,layout_name,page_note=page_note,role=role)
                    i2+=max(req,1)
                    page_idx+=1
                ev_idx+=1

        self.log_cb("PPTX: 封底...")
        self._slide_back(prs,blank,all_photos,album_title)
        prs.save(self.output_path)
        self.log_cb(f"PPTX 已保存: {self.output_path}")

    # ── EMU helpers matching PDF layout exactly ────────────────────────────
    def _mm(self, v): return Emu(int(v*36000))
    def _pt2emu(self, pt): return Emu(int(pt*12700))  # 1pt = 12700 EMU

    # Convert PDF coordinates (points) to PPTX EMU
    # PDF: origin bottom-left; PPTX: origin top-left
    def _x(self, pt_x): return Emu(int(pt_x/PAGE_W * A4_EMU_W))
    def _y(self, pt_y): return Emu(int((PAGE_H-pt_y)/PAGE_H * A4_EMU_H))
    def _w(self, pt_w): return Emu(int(pt_w/PAGE_W * A4_EMU_W))
    def _h(self, pt_h): return Emu(int(pt_h/PAGE_H * A4_EMU_H))

    def _slide_cover(self, prs, blank, all_photos, title, dates, monthly_groups):
        s=prs.slides.add_slide(blank); W,H=self.SW,self.SH
        pal=self._pal(dates[0].month if dates else None)
        best=sorted(all_photos,key=lambda p:p.score,reverse=True)
        # Hero bg
        if best:
            try:
                thumb=self.cache.get_or_make(best[0].path)
                img=ImageEnhance.Brightness(thumb).enhance(0.55)
                self._add_img(s,img,0,0,W,H)
            except: pass
        # Dark overlay
        ov=s.shapes.add_shape(1,0,0,W,H); ov.fill.solid()
        dark=pal.get("cover_dark","#0D1B2A"); r,g,b=[int(c*255) for c in hex_to_rgb(dark)]
        ov.fill.fore_color.rgb=RGBColor(r,g,b); ov.line.fill.background()
        # Left accent stripe
        stripe=s.shapes.add_shape(1,0,0,self._mm(4),H)
        stripe.fill.solid(); stripe.fill.fore_color.rgb=self._rgb(pal["accent"]); stripe.line.fill.background()
        # Title (lower 40% of page)
        ty_emu=Emu(int(A4_EMU_H*0.52))
        self._textbox(s,title,self._mm(MARGIN/mm),ty_emu,W-self._mm(MARGIN/mm*2),self._mm(22),
                      size=36,bold=True,color="FFFFFF",align=PP_ALIGN.CENTER)
        # Accent line
        ln=s.shapes.add_shape(1,self._mm(MARGIN/mm+20),ty_emu+self._mm(24),
                               W-self._mm((MARGIN/mm+20)*2),self._mm(0.5))
        ln.fill.solid(); ln.fill.fore_color.rgb=self._rgb(pal["accent"]); ln.line.fill.background()
        # Date span
        if dates:
            span=f"{dates[0].strftime('%Y.%m.%d')}  -  {dates[-1].strftime('%Y.%m.%d')}"
            self._textbox(s,span,self._mm(MARGIN/mm),ty_emu+self._mm(28),W-self._mm(MARGIN/mm*2),self._mm(10),
                          size=11,color="DDDDDD",align=PP_ALIGN.CENTER)
        # Timeline
        tl_events=[]
        for _y,_m,evts in monthly_groups:
            for grp,loc in evts:
                d=sorted([p.dt for p in grp if p.dt])
                if d: tl_events.append((d[0],loc))
                if len(tl_events)>=7: break
            if len(tl_events)>=7: break
        if tl_events:
            n=len(tl_events); lx0=self._mm(MARGIN/mm+5); lx1=W-self._mm(MARGIN/mm+5)
            line_y=H-self._mm(MARGIN/mm+8)
            # baseline
            ln2=s.shapes.add_shape(1,lx0,line_y,lx1-lx0,self._mm(0.4))
            ln2.fill.solid(); ln2.fill.fore_color.rgb=RGBColor(180,180,180); ln2.line.fill.background()
            spc=(lx1-lx0)//max(n-1,1)
            for i,(dt,loc) in enumerate(tl_events):
                cx=lx0+spc*i
                dot=s.shapes.add_shape(9,cx-self._mm(2),line_y-self._mm(2),self._mm(4),self._mm(4))
                dot.fill.solid(); dot.fill.fore_color.rgb=self._rgb(pal["accent"]); dot.line.fill.background()
                self._textbox(s,dt.strftime("%m/%d"),cx-self._mm(6),line_y-self._mm(9),
                              self._mm(12),self._mm(6),size=7,color="FFFFFF",align=PP_ALIGN.CENTER)

    def _slide_toc(self, prs, blank, monthly_groups):
        s=prs.slides.add_slide(blank); W,H=self.SW,self.SH
        pal=self._pal(monthly_groups[0][1] if monthly_groups else None)
        self._bg_rect(s,pal["bg"])
        # Title
        self._textbox(s,"目  录",self._mm(MARGIN/mm),self._mm(MARGIN/mm),
                      W-self._mm(MARGIN/mm*2),self._mm(16),size=22,bold=True,
                      color=pal["accent"].lstrip("#"))
        # Accent line
        ln=s.shapes.add_shape(1,self._mm(MARGIN/mm),self._mm(MARGIN/mm+18),
                               W-self._mm(MARGIN/mm*2),self._mm(0.5))
        ln.fill.solid(); ln.fill.fore_color.rgb=self._rgb(pal["accent"]); ln.line.fill.background()
        y=self._mm(MARGIN/mm+24)
        row_h=self._mm(7)
        for mi,(year,month,evts) in enumerate(monthly_groups):
            if y>H-self._mm(15): break
            show_month = EventClusterer.should_create_month_page(evts)
            if show_month:
                mlabel=f"{year}年 {MONTH_ZH[month-1]}" if month>0 else "未知"
                bar=s.shapes.add_shape(1,self._mm(MARGIN/mm),y,self._mm(2),self._mm(5))
                bar.fill.solid(); bar.fill.fore_color.rgb=self._rgb(pal["accent"]); bar.line.fill.background()
                self._textbox(s,mlabel,self._mm(MARGIN/mm+4),y,self._mm(60),row_h,
                              size=12,bold=True,color=pal["text"].lstrip("#"))
                y+=self._mm(8)
            for ei,(grp,loc) in enumerate(evts):
                if y>H-self._mm(10): break
                entry=f"  • {loc or 'Unknown'}"
                self._textbox(s,entry,self._mm(MARGIN/mm+2),y,W-self._mm(MARGIN/mm*2+4),row_h,
                              size=10,color=pal["text"].lstrip("#"))
                y+=self._mm(6.5)
            y+=self._mm(3 if show_month else 1.5)

    def _slide_month(self, prs, blank, year, month):
        s=prs.slides.add_slide(blank); W,H=self.SW,self.SH
        pal=self._pal(month); self._bg_rect(s,pal["accent"])
        zh=MONTH_ZH[month-1] if month>0 else "未知"
        en=(MONTH_EN[month-1].upper() if month>0 else "")
        self._textbox(s,en,0,self._mm(50),W,self._mm(60),size=72,bold=True,
                      color="FFFFFF",align=PP_ALIGN.CENTER)
        self._textbox(s,zh,0,self._mm(120),W,self._mm(35),size=42,bold=True,
                      color="FFFFFF",align=PP_ALIGN.CENTER)
        self._textbox(s,str(year),0,self._mm(160),W,self._mm(12),size=14,
                      color="FFFFFF",align=PP_ALIGN.CENTER)

    def _slide_chapter(self, prs, blank, idx, group, location_label, month, story_plan=None):
        s=prs.slides.add_slide(blank); W,H=self.SW,self.SH; pal=self._pal(month)
        self._bg_rect(s,pal["bg"])
        # Left column photos
        cw=Emu(int(A4_EMU_W*0.45)); photos=group[:3]
        if photos:
            sh=H//len(photos)
            for i,ph in enumerate(photos): self._add_pic(s,ph,0,sh*i,cw,sh)
        # Dark overlay on photos
        ov=s.shapes.add_shape(1,0,0,cw,H); ov.fill.solid()
        ov.fill.fore_color.rgb=RGBColor(0,0,0); ov.line.fill.background()
        # Right side info
        rx=cw+self._mm(MARGIN/mm); ry=self._mm(70)
        rw=W-rx-self._mm(MARGIN/mm)
        self._textbox(s,f"{idx+1:02d}",rx,ry,rw,self._mm(30),
                      size=48,bold=True,color=pal["accent"].lstrip("#"))
        ln=s.shapes.add_shape(1,rx,ry+self._mm(34),rw,self._mm(0.5))
        ln.fill.solid(); ln.fill.fore_color.rgb=self._rgb(pal["accent"]); ln.line.fill.background()
        if story_plan:
            self._textbox(s,story_plan.chapter_kicker,rx,ry+self._mm(31),rw,self._mm(6),
                          size=8,color=pal["accent"].lstrip("#"))
        clean_loc = _clean_place_for_display(location_label)
        gps_c=[p for p in group if p.lat is not None]
        has_named_location = bool(clean_loc) and not StoryEngine._is_non_location(clean_loc)
        title_text = clean_loc if has_named_location else ""
        self._textbox(s,title_text,rx,ry+self._mm(37),rw,self._mm(14),
                      size=16,bold=True,color=pal["text"].lstrip("#"))
        dates=sorted([p.dt for p in group if p.dt])
        meta_y = ry + self._mm(54)
        rows = []
        if dates:
            dstr=dates[0].strftime("%Y年%m月%d日") if dates[0].date()==dates[-1].date() \
                 else f"{dates[0].strftime('%m月%d日')} - {dates[-1].strftime('%m月%d日')}"
            rows.append(("日期", dstr))
        if gps_c:
            rep=gps_c[len(gps_c)//2]
            gps=f"{abs(rep.lat):.4f}{'N' if rep.lat>=0 else 'S'}  {abs(rep.lon):.4f}{'E' if rep.lon>=0 else 'W'}"
            rows.append(("坐标", gps))
        if story_plan:
            rows.append(("片段", story_plan.summary_line))
        for label, value in rows:
            self._textbox(s,label,rx,meta_y,self._mm(11),self._mm(6),
                          size=7.2,color=pal["accent"].lstrip("#"))
            self._textbox(s,value,rx+self._mm(12),meta_y-self._mm(0.3),rw-self._mm(12),self._mm(8),
                          size=8.5,color=pal["sub"].lstrip("#"))
            meta_y += self._mm(7.5)
        story_loc = clean_loc if has_named_location else ""
        story = story_plan.chapter_story if story_plan else StoryEngine.generate(group, story_loc, self.cache)
        if story:
            self._textbox(s,story,rx,meta_y+self._mm(2),rw,self._mm(28),
                          size=8,color=pal["sub"].lstrip("#"))
        self._add_portrait_asset(
            s,
            group,
            rx,
            H-self._mm(16),
            int(rw*0.92),
            int(H*0.26),
        )

    def _slide_photos(self, prs, blank, photos, month, layout_name="diptych", page_note="", role="narrative"):
        """Mirror the PDF layout as closely as possible in PPTX."""
        s=prs.slides.add_slide(blank); W,H=self.SW,self.SH; pal=self._pal(month)
        self._bg_rect(s,pal["bg"])
        n=len(photos)
        if n==0: return
        MG=self._mm(MARGIN/mm); IMG_Y0=self._mm((MARGIN+22*mm)/mm)
        IW=W-2*MG; IH=H-IMG_Y0-self._mm(8)
        GAP=self._mm(3)

        if layout_name in ("single","wide_banner","opening","highlight"):
            # One photo full area
            self._add_pic(s,photos[0],MG,IMG_Y0,IW,IH)
        elif layout_name=="diptych" and n>=2:
            w=(IW-GAP)//2
            for i,ph in enumerate(photos[:2]):
                self._add_pic(s,ph,MG+i*(w+GAP),IMG_Y0,w,IH)
        elif layout_name in ("triptych","asymmetric3") and n>=3:
            lw=Emu(int(IW*0.58)); rw=IW-lw-GAP; hh=(IH-GAP)//2
            self._add_pic(s,photos[0],MG,IMG_Y0,lw,IH)
            self._add_pic(s,photos[1],MG+lw+GAP,IMG_Y0+hh+GAP,rw,hh)
            self._add_pic(s,photos[2],MG+lw+GAP,IMG_Y0,rw,hh)
        elif layout_name=="quad" and n>=4:
            w=(IW-GAP)//2; h=(IH-GAP)//2
            pos=[(0,1),(1,1),(0,0),(1,0)]
            for i,ph in enumerate(photos[:4]):
                c2,r2=pos[i]; self._add_pic(s,ph,MG+c2*(w+GAP),IMG_Y0+r2*(h+GAP),w,h)
        elif layout_name=="mosaic5" and n>=5:
            th=Emu(int(IH*0.52)); bh=IH-th-GAP; tw=(IW-GAP)//2; bw=(IW-2*GAP)//3
            for i,ph in enumerate(photos[:2]):
                self._add_pic(s,ph,MG+i*(tw+GAP),IMG_Y0+bh+GAP,tw,th)
            for i,ph in enumerate(photos[2:5]):
                self._add_pic(s,ph,MG+i*(bw+GAP),IMG_Y0,bw,bh)
        elif layout_name in ("hero_thumbs","magazine") and n>=3:
            hh=Emu(int(IH*0.62)); bh=IH-hh-GAP; bw=(IW-GAP)//2
            self._add_pic(s,photos[0],MG,IMG_Y0+bh+GAP,IW,hh)
            for j,ph in enumerate(photos[1:3]):
                self._add_pic(s,ph,MG+j*(bw+GAP),IMG_Y0,bw,bh)
        elif layout_name=="strip" and n>=3:
            sw=(IW-GAP*(n-1))//n
            for i,ph in enumerate(photos[:n]):
                self._add_pic(s,ph,MG+i*(sw+GAP),IMG_Y0,sw,IH)
        elif layout_name=="circle_cluster" and n>=1:
            gap2=self._mm(4)
            big_d=min(Emu(int(IW*0.48)), Emu(int(IH*0.78)))
            remaining=IW-big_d-gap2
            small_d=min(Emu(int((remaining-gap2)/2)), Emu(int((IH-gap2)/2)))
            slots=[
                (MG+self._mm(2), IMG_Y0+(IH-big_d)//2, big_d, big_d, 0),
                (MG+big_d+gap2, IMG_Y0+IH-small_d, small_d, small_d, 1),
                (MG+big_d+gap2+small_d+gap2, IMG_Y0+IH-small_d, small_d, small_d, 2),
                (MG+big_d+gap2, IMG_Y0, small_d, small_d, 3),
                (MG+big_d+gap2+small_d+gap2, IMG_Y0, small_d, small_d, 4),
            ]
            for left,top,w2,h2,idx in slots:
                if idx < n:
                    self._add_pic_masked(s,photos[idx],left,top,w2,h2,shape="circle")
        elif layout_name=="hexagon_hive" and n>=1:
            gap2=self._mm(2.6)
            big_w=min(Emu(int(IW*0.33)), Emu(int(IH*0.34))); big_h=Emu(int(big_w*0.88))
            small_w=Emu(int(big_w*0.74)); small_h=Emu(int(small_w*0.88))
            cx=MG+IW//2; cy=IMG_Y0+IH//2
            x_radius=Emu(int(big_w*0.78 + small_w*0.56 + gap2))
            y_radius=Emu(int(big_h*0.86 + small_h*0.56 + gap2))
            diag_x=Emu(int(x_radius*0.88))
            diag_y=Emu(int(y_radius*0.52))

            def centered(center_x, center_y, w2, h2):
                return center_x-w2//2, center_y-h2//2, w2, h2

            slots=[
                centered(cx, cy, big_w, big_h),
                centered(cx, cy+y_radius, small_w, small_h),
                centered(cx+diag_x, cy+diag_y, small_w, small_h),
                centered(cx+diag_x, cy-diag_y, small_w, small_h),
                centered(cx, cy-y_radius, small_w, small_h),
                centered(cx-diag_x, cy-diag_y, small_w, small_h),
                centered(cx-diag_x, cy+diag_y, small_w, small_h),
            ]
            for idx,(left,top,w2,h2) in enumerate(slots):
                if idx < n:
                    self._add_pic_masked(s,photos[idx],left,top,w2,h2,shape="hexagon")
        elif layout_name=="nine_grid" and n>=1:
            cols=rows=3
            gap2=self._mm(2.2)
            cw2=(IW-gap2*(cols-1))//cols; ch2=(IH-gap2*(rows-1))//rows
            for i,ph in enumerate(photos[:9]):
                c2=i%cols; r2=i//cols
                self._add_pic(s,ph,MG+c2*(cw2+gap2),IMG_Y0+r2*(ch2+gap2),cw2,ch2)
        else:
            # Fallback: grid
            cols=min(n,2); rows=math.ceil(n/cols)
            cw2=(IW-GAP*(cols-1))//cols; ch2=(IH-GAP*(rows-1))//rows
            for i,ph in enumerate(photos):
                c2=i%cols; r2=i//cols
                self._add_pic(s,ph,MG+c2*(cw2+GAP),IMG_Y0+r2*(ch2+GAP),cw2,ch2)

        # Caption strip
        ph=photos[0]; parts=[]
        if ph.dt: parts.append(ph.date_str+(" "+ph.time_str if ph.time_str else ""))
        if ph.raw_gps_str: parts.append(ph.raw_gps_str)
        cap=" | ".join(parts)
        cap_y=H-self._mm(7)
        self._textbox(s,cap,MG,cap_y,IW,self._mm(6),size=7,color=pal["sub"].lstrip("#"))
        self._slide_story_note(s, page_note, role, month)

    def _slide_back(self, prs, blank, all_photos, album_title):
        s=prs.slides.add_slide(blank); W,H=self.SW,self.SH; pal=self._pal(None)
        best=sorted(all_photos,key=lambda p:p.score,reverse=True)[:12]
        if len(best)<4: best=best*4
        cols,rows=4,3; cw=W//cols; rh=H//rows
        random.seed(99)
        for i,ph in enumerate(best[:12]):
            col=i%cols; row=i//cols; self._add_pic(s,ph,col*cw,row*rh,cw,rh)
        # No tint overlay — photos shown at original colour
        bw=self._mm(120); bh=self._mm(20); bx=(W-bw)//2; by=(H-bh)//2
        badge=s.shapes.add_shape(1,bx,by,bw,bh)
        badge.fill.solid(); badge.fill.fore_color.rgb=RGBColor(0,0,0); badge.line.fill.background()
        self._textbox(s,album_title,bx,by,bw,bh,size=28,bold=True,
                      color="FFFFFF",align=PP_ALIGN.CENTER)

# ══ MODULE 6: Orchestrator (pipelined, threaded) ════════════════════════════

def scan_exif_only(folder_list, log_cb=None, exclude_dirs=None):
    """
    Stage 1: Fast EXIF-only scan. No pixel loading.
    Returns list[PhotoInfo] sorted by datetime.
    """
    log=log_cb or default_log; all_files=[]
    exclude_roots = {
        str(Path(path).resolve()).lower()
        for path in (exclude_dirs or [])
        if path
    }

    def should_skip_dir(path: str) -> bool:
        resolved = str(Path(path).resolve()).lower()
        name = Path(path).name.lower()
        if name in {".album_cache", "__pycache__", ".codex_pycache"}:
            return True
        return any(
            resolved == root or resolved.startswith(root + os.sep)
            for root in exclude_roots
        )

    for folder in folder_list:
        for root, dirs, files in os.walk(folder):
            dirs[:] = [d for d in dirs if not should_skip_dir(os.path.join(root, d))]
            for f in files:
                if Path(f).suffix.lower() in SUPPORTED_EXTS:
                    all_files.append(os.path.join(root,f))
    log(f"扫描到 {len(all_files)} 张照片，读取 EXIF 元数据...")
    photos=[]
    for fp in all_files:
        try: photos.append(PhotoInfo(fp))
        except: pass
    return photos

def score_candidates(photos, cache, max_workers=4, log_cb=None, progress_cb=None):
    """Stage 3: Score only candidate photos (already selected per event)."""
    log=log_cb or default_log; log(f"多线程评分 {len(photos)} 张候选照片（{max_workers} 线程）...")
    PhotoScorer.score_batch(photos, cache, max_workers=max_workers, progress_cb=progress_cb)

def build_album(folder_list, pdf_path, pptx_path,
                album_title="我的相册",
                palette_name="四季自动",
                time_gap_hours=6.0,
                geo_radius_km=30.0,
                max_per_event=30,
                score_photos=True,
                resolve_geo=True,
                export_pptx=True,
                max_workers=4,
                cache_dir=None,
                portrait_assets_dir=None,
                enable_portrait_elements=False,
                progress_cb=None,
                log_cb=None):

    log  = log_cb or default_log
    prog = progress_cb or (lambda v:None)
    t0   = time.time()

    # ── Stage 1: EXIF scan (no image loading) ──────────────────────────────
    # ── Stage 0: Plan exclusions / cache ──────────────────────────────────
    if cache_dir is None:
        cache_dir = Path(folder_list[0]) / ".album_cache"
    source_root_paths = [Path(folder).resolve() for folder in folder_list]
    source_roots = {str(root).lower() for root in source_root_paths}
    exclude_dirs = [cache_dir]
    for output_path in [pdf_path, pptx_path]:
        if not output_path:
            continue
        parent = Path(output_path).resolve().parent
        parent_str = str(parent).lower()
        is_source_root = parent_str in source_roots
        is_under_source_root = any(
            parent == root or root in parent.parents
            for root in source_root_paths
        )
        if not is_source_root and (is_under_source_root or parent_str not in source_roots):
            exclude_dirs.append(parent)

    # ── Stage 1: EXIF scan (no image loading) ──────────────────────────────
    photos = scan_exif_only(folder_list, log, exclude_dirs=exclude_dirs)
    if not photos: raise ValueError("所选文件夹中未找到支持的照片格式。")
    prog(5)

    # ── Stage 2: Cluster by time+geo using only EXIF data ─────────────────
    log(f"事件聚类（间隔 {time_gap_hours:.0f}h，半径 {geo_radius_km:.0f}km）...")
    clusterer = EventClusterer(time_gap_hours, geo_radius_km)
    groups    = clusterer.cluster(photos)
    raw_event_count = len(groups)
    groups = EventClusterer.merge_small_events(groups, min_photos=EventClusterer.MIN_EVENT_PHOTOS)
    log(f"  → 初始 {raw_event_count} 个事件，合并小事件后 {len(groups)} 个事件，共 {len(photos)} 张")

    prog(10)

    # ── Stage 4: Setup thumbnail cache ────────────────────────────────────
    cache = ThumbCache(cache_dir, thumb_size=1024)
    log(f"缩略图缓存: {cache_dir}")
    init_geocache(cache_dir)   # load geocache.json if it exists

    portrait_store = None
    portrait_assets: list[PortraitAsset] = []
    if enable_portrait_elements:
        portrait_store = PortraitAssetStore(
            portrait_assets_dir or get_default_portrait_library_path()
        )
        portrait_assets = portrait_store.list_assets()
        if portrait_assets:
            log(f"人像元素素材库: {portrait_store.root_dir}（{len(portrait_assets)} 个 PNG 元素）")
        else:
            log(f"人像元素素材库为空，已跳过叠加: {portrait_store.root_dir}")
            portrait_store = None

    # ── Stage 5: Parallel scoring (sample only for speed, ALL photos kept) ──
    # IMPORTANT: max_per_event limits how many photos get SCORED per event,
    # but ALL photos are retained in the album — nothing is deleted.
    # For 5000 photos this means we score ~300 samples instead of all 5000,
    # saving significant time. Unscored photos get the default score=0.3.
    if score_photos:
        score_sample = [p for g in groups for p in g[:max_per_event]]
        total_photos = sum(len(g) for g in groups)
        log(f"多线程评分（样本 {len(score_sample)} 张 / 全部 {total_photos} 张，"
            f"{max_workers} 线程）。所有照片均保留在相册中。")
        def score_prog(v): prog(10+int(v*0.20))
        score_candidates(score_sample, cache, max_workers=max_workers,
                        log_cb=log, progress_cb=score_prog)
        # Sort each group by datetime (preserve chronological order for layout)
        groups = [sorted(g, key=lambda p: p.dt or datetime.min) for g in groups]
    prog(30)

    # ── Stage 6: Reverse geocoding — ONE query per EVENT only ───────────
    # Only geocodes the representative photo of each event, not every photo.
    # 20 events → at most 20 network calls. Results cached in geocache.json.
    # ph.location_name is NEVER set — place names go into location_names[] only.
    group_place_names: list[str] = [""] * len(groups)
    if resolve_geo:
        log("📍 事件代表坐标反向地理编码（每事件1次，带磁盘缓存）...")
        # Pick one representative GPS coordinate per group (median photo)
        rep_coords: list[tuple | None] = []
        for grp in groups:
            gps_cands = [p for p in grp if p.lat is not None]
            if gps_cands:
                rep = gps_cands[len(gps_cands)//2]
                rep_coords.append((rep.lat, rep.lon))
            else:
                rep_coords.append(None)
        n_needed = sum(1 for c in rep_coords if c is not None)
        log(f"  → {n_needed} 个事件需要查询（共 {len(groups)} 个事件）")
        # Resolve unique coords in parallel (fills _geocache as side-effect)
        unique = list({latlng for latlng in rep_coords if latlng is not None})
        with ThreadPoolExecutor(max_workers=min(4, max_workers)) as ex:
            list(ex.map(lambda ll: reverse_geocode(ll[0], ll[1]), unique))
        # Build per-group place name list — do NOT set ph.location_name
        for i, latlng in enumerate(rep_coords):
            group_place_names[i] = reverse_geocode(latlng[0], latlng[1]) if latlng else ""
        log("  → 地理编码完成")

    # Derive final location label per group:
    #   1. Geocoded English place name  (e.g. "El Paso, Texas, United States")
    #   2. Raw GPS string               (e.g. "31.7667N  106.5031W")
    #   3. Date string                  (last resort)
    location_names: list[str] = []
    for grp, place in zip(groups, group_place_names):
        if place:
            location_names.append(place)
        else:
            gps_cands = [p for p in grp if p.lat is not None]
            if gps_cands:
                rep = gps_cands[len(gps_cands)//2]
                location_names.append(gps_str(rep.lat, rep.lon))
            else:
                dates2 = sorted([p.dt for p in grp if p.dt])
                location_names.append(
                    dates2[0].strftime("%Y年%m月%d日") if dates2 else "旅途中")

    elapsed = time.time()-t0
    log(f"预处理完成，耗时 {elapsed:.1f}s")

    # ── Stage 7: Build PDF ─────────────────────────────────────────────────
    log("生成 PDF 相册...")
    pdf_builder = AlbumBuilder(pdf_path, palette_name, cache,
                               progress_cb=lambda v: prog(30+int(v*0.55)),
                               log_cb=log,
                               portrait_store=portrait_store,
                               portrait_assets=portrait_assets)
    pdf_builder.build(groups, location_names, album_title)

    # ── Stage 8: Build PPTX ────────────────────────────────────────────────
    if export_pptx and pptx_path:
        if not HAS_PPTX:
            log("python-pptx 未安装，跳过 PPTX。pip install python-pptx")
        else:
            try:
                log("生成 PPTX (A4)...")
                PptxBuilder(
                    pptx_path,
                    palette_name,
                    cache,
                    log_cb=log,
                    portrait_store=portrait_store,
                    portrait_assets=portrait_assets,
                )\
                    .build(groups, location_names, album_title)
            except Exception as e:
                import traceback
                log(f"PPTX 失败（PDF 正常）: {e}\n{traceback.format_exc()}")
    prog(100)
    log(f"全部完成，总耗时 {time.time()-t0:.1f}s")


def build_album_from_request(request: AlbumBuildRequest):
    return build_album(
        folder_list=request.folder_list,
        pdf_path=request.pdf_path,
        pptx_path=request.pptx_path,
        album_title=request.album_title,
        palette_name=request.palette_name,
        time_gap_hours=request.time_gap_hours,
        geo_radius_km=request.geo_radius_km,
        max_per_event=request.max_per_event,
        score_photos=request.score_photos,
        resolve_geo=request.resolve_geo,
        export_pptx=request.export_pptx,
        max_workers=request.max_workers,
        cache_dir=request.cache_dir,
        portrait_assets_dir=request.portrait_assets_dir,
        enable_portrait_elements=request.enable_portrait_elements,
        progress_cb=request.progress_cb,
        log_cb=request.log_cb,
    )


# ══ MODULE 7: GUI ════════════════════════════════════════════════════════════

class App(tk.Tk):
    def __init__(self):
        super().__init__(); self.title("智能照片相册生成器 v4.0")
        self.resizable(False,False); self.configure(bg="#F0F4F8")
        self.folder_list=[]; self._build_ui()

    def _build_ui(self):
        BG="#F0F4F8"; ACC="#2563EB"
        sty=ttk.Style(self); sty.theme_use("clam")
        sty.configure("TButton",padding=6,relief="flat",background=ACC,foreground="white",font=("Helvetica",10))
        sty.map("TButton",background=[("active","#1D4ED8")])
        sty.configure("TLabel",background=BG,font=("Helvetica",10))
        sty.configure("TEntry",padding=4); sty.configure("TCombobox",padding=4)

        hdr=tk.Frame(self,bg=ACC,height=60); hdr.pack(fill="x")
        tk.Label(hdr,text="📷 智能照片相册生成器 v4.0",
                 font=("Helvetica",16,"bold"),bg=ACC,fg="white").pack(pady=14)

        body=tk.Frame(self,bg=BG); body.pack(fill="both",expand=True,padx=20,pady=12)

        tk.Label(body,text="📁 照片文件夹",font=("Helvetica",11,"bold"),bg=BG).grid(
            row=0,column=0,columnspan=3,sticky="w",pady=(0,4))
        ff=tk.Frame(body,bg=BG); ff.grid(row=1,column=0,columnspan=3,sticky="ew")
        self.folder_lb=tk.Listbox(ff,height=4,width=62,font=("Helvetica",9),
                                   selectmode=tk.EXTENDED,relief="solid",bd=1,highlightthickness=0)
        sc=ttk.Scrollbar(ff,orient="vertical",command=self.folder_lb.yview)
        self.folder_lb.configure(yscrollcommand=sc.set)
        self.folder_lb.pack(side="left",fill="both",expand=True); sc.pack(side="right",fill="y")
        bf=tk.Frame(body,bg=BG); bf.grid(row=2,column=0,columnspan=3,sticky="w",pady=4)
        ttk.Button(bf,text="+ 添加文件夹",command=self._add_folders).pack(side="left",padx=(0,6))
        ttk.Button(bf,text="x 移除选中",command=self._remove_folders).pack(side="left")

        ttk.Separator(body,orient="horizontal").grid(row=3,column=0,columnspan=3,sticky="ew",pady=8)

        def lbl(r,t): tk.Label(body,text=t,bg=BG).grid(row=r,column=0,sticky="w")
        def entry(r,var,w=12):
            ttk.Entry(body,textvariable=var,width=w).grid(row=r,column=1,sticky="w",padx=6)

        lbl(4,"📝 相册标题"); self.title_var=tk.StringVar(value="我的相册")
        ttk.Entry(body,textvariable=self.title_var,width=30).grid(row=4,column=1,sticky="ew",padx=6)

        lbl(5,"🎨 配色主题"); self.palette_var=tk.StringVar(value="四季自动")
        ttk.Combobox(body,textvariable=self.palette_var,
                     values=list(MANUAL_PALETTES.keys()),state="readonly",width=28).grid(
            row=5,column=1,sticky="ew",padx=6,pady=4)

        lbl(6,"⏱ 事件间隔（小时）"); self.gap_var=tk.StringVar(value="6"); entry(6,self.gap_var)
        lbl(7,"📍 地理半径（千米）"); self.geo_km_var=tk.StringVar(value="30"); entry(7,self.geo_km_var)
        lbl(8,"🔢 每事件最多张数"); self.max_var=tk.StringVar(value="30"); entry(8,self.max_var)
        lbl(9,"⚡ 评分线程数"); self.threads_var=tk.StringVar(value="4"); entry(9,self.threads_var)

        self.score_var=tk.BooleanVar(value=True)
        ttk.Checkbutton(body,text="🔍 照片质量评分（不删除）",variable=self.score_var).grid(
            row=10,column=0,columnspan=2,sticky="w",pady=(6,2))
        self.geo_var=tk.BooleanVar(value=True)
        ttk.Checkbutton(body,text="🌍 在线解析地理位置（带缓存，重复运行极快）",
                        variable=self.geo_var).grid(row=10,column=0,columnspan=2,sticky="w",pady=(6,2))
        self.pptx_var=tk.BooleanVar(value=True)
        ttk.Checkbutton(body,text="📊 生成 A4 PPTX（需 pip install python-pptx）",
                        variable=self.pptx_var).grid(row=11,column=0,columnspan=2,sticky="w")
        self.portrait_var=tk.BooleanVar(value=True)
        ttk.Checkbutton(body,text="🧩 在章节页加入人像去背元素（来自素材库）",
                        variable=self.portrait_var).grid(row=12,column=0,columnspan=2,sticky="w")

        ttk.Separator(body,orient="horizontal").grid(row=13,column=0,columnspan=3,sticky="ew",pady=8)

        lbl(14,"💾 PDF 输出"); self.pdf_var=tk.StringVar(value=str(Path.home()/"photo_album.pdf"))
        ttk.Entry(body,textvariable=self.pdf_var,width=38).grid(row=14,column=1,sticky="ew",padx=6)
        ttk.Button(body,text="浏览",command=self._pick_pdf).grid(row=14,column=2)

        lbl(15,"📊 PPTX 输出"); self.pptx_out_var=tk.StringVar(value=str(Path.home()/"photo_album.pptx"))
        ttk.Entry(body,textvariable=self.pptx_out_var,width=38).grid(row=15,column=1,sticky="ew",padx=6,pady=4)
        ttk.Button(body,text="浏览",command=self._pick_pptx).grid(row=15,column=2)

        lbl(16,"📁 缩略图缓存目录"); self.cache_var=tk.StringVar(value="")
        ttk.Entry(body,textvariable=self.cache_var,width=38).grid(row=16,column=1,sticky="ew",padx=6)
        ttk.Button(body,text="浏览",command=self._pick_cache).grid(row=16,column=2)
        lbl(17,"🧩 人像素材目录"); self.portrait_dir_var=tk.StringVar(value=str(get_default_portrait_library_path()))
        ttk.Entry(body,textvariable=self.portrait_dir_var,width=38).grid(row=17,column=1,sticky="ew",padx=6,pady=4)
        ttk.Button(body,text="浏览",command=self._pick_portrait_dir).grid(row=17,column=2)
        tk.Label(body,text="缓存目录留空则自动创建在照片目录下",bg=BG,fg="#888",
                 font=("Helvetica",8)).grid(row=18,column=1,sticky="w",padx=6)
        tk.Label(body,text="去背工具的“保存到相册素材库”默认会写入这里；留空则使用默认 portrait_elements",bg=BG,fg="#888",
                 font=("Helvetica",8)).grid(row=19,column=1,sticky="w",padx=6)

        self.progress=ttk.Progressbar(body,orient="horizontal",length=500,mode="determinate")
        self.progress.grid(row=20,column=0,columnspan=3,pady=10,sticky="ew")
        self.log_text=tk.Text(body,height=8,width=62,state="disabled",
                               font=("Courier",8),bg="#1E1E2E",fg="#CDD6F4",
                               relief="solid",bd=1,highlightthickness=0)
        self.log_text.grid(row=21,column=0,columnspan=3,sticky="ew")
        self.gen_btn=ttk.Button(body,text="🚀  生成相册",command=self._start_build)
        self.gen_btn.grid(row=22,column=0,columnspan=3,pady=12,ipadx=20,ipady=4)

        warnings=[]
        if not HEIC_SUPPORTED: warnings.append("⚠ 未安装 pillow-heif，HEIC 不支持")
        if not HAS_PPTX:       warnings.append("⚠ 未安装 python-pptx，PPTX 不可用")
        for i,w2 in enumerate(warnings):
            tk.Label(body,text=w2,bg=BG,fg="#B45309",font=("Helvetica",8)).grid(
                row=23+i,column=0,columnspan=3)

    def _add_folders(self):
        d=filedialog.askdirectory(mustexist=True,title="选择照片文件夹")
        if d and d not in self.folder_list: self.folder_list.append(d); self.folder_lb.insert(tk.END,d)
    def _remove_folders(self):
        for i in reversed(list(self.folder_lb.curselection())): self.folder_lb.delete(i); del self.folder_list[i]
    def _pick_pdf(self):
        p=filedialog.asksaveasfilename(defaultextension=".pdf",filetypes=[("PDF","*.pdf")],initialfile="photo_album.pdf")
        if p: self.pdf_var.set(p)
    def _pick_pptx(self):
        p=filedialog.asksaveasfilename(defaultextension=".pptx",filetypes=[("PPTX","*.pptx")],initialfile="photo_album.pptx")
        if p: self.pptx_out_var.set(p)
    def _pick_cache(self):
        d=filedialog.askdirectory(title="选择缓存目录")
        if d: self.cache_var.set(d)
    def _pick_portrait_dir(self):
        d=filedialog.askdirectory(title="选择人像素材目录")
        if d: self.portrait_dir_var.set(d)
    def _log(self,msg):
        def _do():
            self.log_text.configure(state="normal"); self.log_text.insert(tk.END,msg+"\n")
            self.log_text.see(tk.END); self.log_text.configure(state="disabled")
        self.after(0,_do)
    def _set_progress(self,v): self.after(0,lambda:self.progress.configure(value=v))
    def _start_build(self):
        if not self.folder_list: messagebox.showwarning("提示","请先添加照片文件夹！"); return
        self.gen_btn.configure(state="disabled"); self.progress.configure(value=0)
        threading.Thread(target=self._worker,daemon=True).start()
    def _worker(self):
        try:
            cache_dir=self.cache_var.get().strip() or None
            request = AlbumBuildRequest(
                folder_list=list(self.folder_list),
                pdf_path=self.pdf_var.get(),
                pptx_path=self.pptx_out_var.get(),
                album_title=self.title_var.get() or "我的相册",
                palette_name=self.palette_var.get(),
                time_gap_hours=float(self.gap_var.get() or 6),
                geo_radius_km=float(self.geo_km_var.get() or 30),
                max_per_event=int(self.max_var.get() or 30),
                score_photos=self.score_var.get(),
                resolve_geo=self.geo_var.get(),
                export_pptx=self.pptx_var.get(),
                max_workers=max(1,int(self.threads_var.get() or 4)),
                cache_dir=cache_dir,
                portrait_assets_dir=self.portrait_dir_var.get().strip() or None,
                enable_portrait_elements=self.portrait_var.get(),
                progress_cb=self._set_progress,
                log_cb=self._log,
            )
            build_album_from_request(request)
            msg=f"PDF: {self.pdf_var.get()}"
            if self.pptx_var.get(): msg+=f"\nPPTX: {self.pptx_out_var.get()}"
            self.after(0,lambda: messagebox.showinfo("完成",f"相册已生成！\n\n{msg}"))
        except Exception as e:
            import traceback; self._log(f"ERROR: {e}\n{traceback.format_exc()}")
            self.after(0,lambda: messagebox.showerror("错误",str(e)))
        finally: self.after(0,lambda: self.gen_btn.configure(state="normal"))

if __name__=="__main__":
    App().mainloop()
